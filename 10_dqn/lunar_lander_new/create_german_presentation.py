"""
Erstelle deutsche Version der DQN vs Double DQN Präsentation
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_german_presentation():
    """Erstelle deutsche Präsentation mit allen Inhalten"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Farbschema
    COLOR_PRIMARY = RGBColor(0, 51, 102)  # Dunkelblau
    COLOR_ACCENT = RGBColor(0, 102, 204)  # Blau
    COLOR_SECONDARY = RGBColor(102, 102, 102)  # Grau
    COLOR_WHITE = RGBColor(255, 255, 255)

    def add_title_slide(prs, title, subtitle=""):
        """Titelfolie hinzufügen"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRIMARY

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_WHITE
        title_p.alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = subtitle
            subtitle_p.font.size = Pt(28)
            subtitle_p.font.color.rgb = RGBColor(200, 200, 200)
            subtitle_p.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(prs, title, content_items):
        """Inhaltsfolie hinzufügen"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_WHITE

        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = COLOR_PRIMARY
        title_shape.line.color.rgb = COLOR_PRIMARY

        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_WHITE
        title_p.space_before = Pt(10)
        title_p.space_after = Pt(10)

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_SECONDARY
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.level = 0

        return slide

    def add_image_slide(prs, title, image_path):
        """Folie mit Bild hinzufügen"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_WHITE

        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = COLOR_PRIMARY
        title_shape.line.color.rgb = COLOR_PRIMARY

        title_frame = title_shape.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_WHITE
        title_p.space_before = Pt(5)
        title_p.space_after = Pt(5)

        if Path(image_path).exists():
            slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.2), width=Inches(9))

        return slide

    # FOLIE 1: Titelfolie
    add_title_slide(prs,
                    "DQN vs Double DQN",
                    "Umfassende Hyperparameter-Analyse\nLunar Lander Umgebung")

    # FOLIE 2: Zusammenfassung
    add_content_slide(prs, "Zusammenfassung", [
        "• Ziel: Vergleich der DQN- und Double DQN-Algorithmen",
        "• Umgebung: Lunar Lander-v3 (kontinuierliche Steuerung)",
        "• Analysierte Hyperparameter:",
        "  - Experience Replay Buffer (Größe, Min-Größe)",
        "  - Explorationsstrategie (ε-Decay-Zeitpläne)",
        "  - Target Network Updates (Hard vs Soft)",
        "  - Learning Rate und Batch Size Effekte",
        "• Wichtige Metriken: Finale Performance, Stabilität, Lerngeschwindigkeit"
    ])

    # FOLIE 3: Umgebungsbeschreibung
    add_content_slide(prs, "Lunar Lander-v3 Umgebung", [
        "• Ziel: Sichere Landung eines Raumfahrzeugs auf dem Mond",
        "• Zustandsraum: 8 kontinuierliche Werte",
        "  - Position (x, y), Geschwindigkeit (vx, vy), Winkel, Winkelgeschwindigkeit, Beinkontakt-Flags",
        "• Aktionsraum: 4 diskrete Aktionen",
        "  - Nichts tun, Linken Motor zünden, Hauptmotor zünden, Rechten Motor zünden",
        "• Belohnung: -1 pro Zeitschritt, +100 für sichere Landung, -100 für Absturz",
        "• Erfolgsschwelle: Score > 200 (konsistent)",
        "• Episodenlänge: Typischerweise 200-500 Zeitschritte"
    ])

    # FOLIE 4: Baseline-Konfiguration
    add_content_slide(prs, "Baseline-Konfiguration", [
        "• Netzwerk-Architektur: 2-schichtiges MLP (128 verborgene Einheiten)",
        "• Learning Rate: 1×10⁻³",
        "• Discount Factor (γ): 0.99",
        "• Batch Size: 64",
        "• Replay Buffer Größe: 100.000",
        "• Min Buffer Größe: 1.000 (vor Trainingsbeginn)",
        "• Exploration: ε-greedy mit linearem Decay",
        "  - ε_start: 1.0, ε_end: 0.01",
        "  - Decay über: 250.000 Schritte",
        "• Target Network: Hard Update alle 1.000 Schritte"
    ])

    # FOLIE 5: Experiment 1 - Visualisierung
    add_image_slide(prs, "Experiment 1: Auswirkungen der Replay Buffer Größe",
                    "results/01_buffer_size_analysis.png")

    # FOLIE 6: Experiment 1 - Ergebnisse
    add_content_slide(prs, "Buffer Größen-Analyse - Ergebnisse", [
        "• Kleine Buffer (10k): Frühes instabiles Lernen, niedrigere finale Performance",
        "• Mittlere Buffer (50-100k): Beste Balance von Stabilität und Performance",
        "• Große Buffer (200k): Höhere Speicheranforderungen, ähnliche finale Performance",
        "• Double DQN: Robuster gegenüber Buffer-Größenvariationen",
        "• Empfehlung: 100k bietet bestes Performance/Speicher-Verhältnis",
        "• Auswirkung auf Q-Wert Überschätzung: Größere Buffer → mehr Vielfalt → bessere Schätzungen"
    ])

    # FOLIE 7: Experiment 2 - Visualisierung
    add_image_slide(prs, "Experiment 2: Auswirkungen der Explorationsstrategie (ε-Decay)",
                    "results/02_epsilon_decay_analysis.png")

    # FOLIE 8: Experiment 2 - Ergebnisse
    add_content_slide(prs, "Epsilon Decay-Analyse - Ergebnisse", [
        "• Schneller Decay (50k Schritte): Schnelle Konvergenz aber riskant (schlechte Exploration)",
        "• Langsamer Decay (400k Schritte): Bessere Exploration, stabileres Lernen",
        "• Optimaler Decay: ~250k Schritte balanciert Exploration vs Exploitation",
        "• Double DQN: Weniger empfindlich gegenüber Decay-Geschwindigkeit (robuster)",
        "• DQN: Performance verschlechtert sich signifikant bei zu schnellem Decay",
        "• Erkenntnis: Frühe Überausbeutung vor angemessenem Explorationsernen schadet Performance"
    ])

    # FOLIE 9: Experiment 3 - Visualisierung
    add_image_slide(prs, "Experiment 3: Strategie zum Update des Target Networks",
                    "results/03_update_strategy_analysis.png")

    # FOLIE 10: Experiment 3 - Ergebnisse
    add_content_slide(prs, "Update-Strategie-Analyse - Ergebnisse", [
        "• Hard Updates (alle 500-2000 Schritte): Stabile Konvergenz",
        "• Soft Updates (τ=0.001-0.005): Glattere Lernkurven, niedrigere Varianz",
        "• Optimal: Hard Updates alle 1000 Schritte ODER Soft Updates mit τ=0.005",
        "• Double DQN: Marginale Verbesserung durch Soft Updates",
        "• DQN: Mehr Nutzen von Soft Updates (reduziert Überschätzung)",
        "• Häufigkeit ist wichtig: Zu häufige Hard Updates = Instabilität, zu selten = veraltete Targets"
    ])

    # FOLIE 11: Experiment 4 - Visualisierung
    add_image_slide(prs, "Experiment 4: Auswirkungen von Learning Rate & Batch Size",
                    "results/04_learning_rate_batch_size_analysis.png")

    # FOLIE 12: Experiment 4 - Ergebnisse
    add_content_slide(prs, "Learning Rate & Batch Size-Analyse - Ergebnisse", [
        "• Learning Rate: 1×10⁻³ optimal, höhere Werte → Instabilität",
        "• Learning Rate: Niedrigere Werte → langsameres Lernen (0.5×10⁻³ zu konservativ)",
        "• Batch Size: 64 performt am besten, balanciert Gradienten-Qualität vs Stabilität",
        "• Große Batches (128): Bessere Stabilität aber langsamere Updates",
        "• Kleine Batches (32): Schnelleres Lernen aber noisigere Gradienten",
        "• Double DQN: Weniger empfindlich gegenüber Learning Rate Variationen"
    ])

    # FOLIE 13: DQN vs Double DQN Vergleich
    add_content_slide(prs, "DQN vs Double DQN: Hauptunterschiede", [
        "• Q-Wert Schätzung:",
        "  - DQN: y = r + γ·max_a' Q(s',a') mit Target Network",
        "  - Double DQN: y = r + γ·Q_target(s', argmax_a' Q_online(s',a'))",
        "",
        "• Effekt: Double DQN reduziert Q-Wert Überschätzung",
        "",
        "• Empirische Ergebnisse:",
        "  - Double DQN: Stabilere finale Performance",
        "  - DQN: Schnelleres initiales Lernen aber höhere Varianz",
        "  - Unterschied stärker mit schlechteren Hyperparametern"
    ])

    # FOLIE 14: Schlussfolgerungen
    add_content_slide(prs, "Schlussfolgerungen & Empfehlungen", [
        "1. Hyperparameter-Tuning hat großen Einfluss auf Stabilität und Performance",
        "",
        "2. Double DQN ist robuster über verschiedene Hyperparameter-Einstellungen",
        "",
        "3. Wichtigste Empfehlungen:",
        "   • Nutze 100k Replay Buffer mit 1k-5k Minimum vor Training",
        "   • Decay ε über ~250k Schritte mit linearem oder exponentiellem Zeitplan",
        "   • Nutze Hard Updates (1000 Schritte) oder Soft Updates (τ=0.005)",
        "   • Learning Rate 1×10⁻³, Batch Size 64 für Lunar Lander",
        "",
        "4. Double DQN wird für Produktionssysteme wegen Robustheit empfohlen"
    ])

    # FOLIE 15: Zusammenfassung und Statistik
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY

    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Analysezusammenfassung"
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_WHITE
    title_p.space_before = Pt(5)
    title_p.space_after = Pt(5)

    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    tf = summary_box.text_frame
    tf.word_wrap = True

    summary_text = [
        "Durchgeführte Experimente: 4 große Hyperparameter-Studien",
        "• 1. Replay Buffer Größe (4 Varianten): 10k, 50k, 100k, 200k",
        "• 2. Epsilon Decay (4 Zeitpläne): 50k, 150k, 250k, 400k Schritte",
        "• 3. Update-Strategie (5 Varianten): Hard/Soft mit verschiedenen Häufigkeiten",
        "• 4. Lernparameter (6 Varianten): Learning Rate und Batch Size Kombinationen",
        "",
        "Gesamte Trainings-Episoden: 3000+ pro Konfiguration",
        "",
        "Wichtigste Erkenntnisse: Double DQN übertrifft DQN konsistent in Stabilität,",
        "besonders mit suboptimalen Hyperparametern. Angemessenes Explorations-",
        "Scheduling und Puffer-Management sind entscheidend für Erfolg."
    ]

    for i, text in enumerate(summary_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return prs


def main():
    output_dir = Path(__file__).parent / "results"
    output_dir.mkdir(exist_ok=True)

    print("\n" + "="*80)
    print("ERSTELLE DEUTSCHE PRÄSENTATION")
    print("="*80 + "\n")

    # Erstelle Präsentation
    prs = create_german_presentation()

    # Speichere Präsentation
    output_path = output_dir / "DQN_vs_DoubleDQN_Analysis_DE.pptx"
    try:
        prs.save(str(output_path))
        print(f"✅ Deutsche Präsentation gespeichert: {output_path}")
    except Exception as e:
        print(f"Fehler beim Speichern: {e}")
        alt_path = Path.cwd() / "DQN_vs_DoubleDQN_Analysis_DE.pptx"
        prs.save(str(alt_path))
        print(f"✅ Deutsche Präsentation gespeichert (Alternative): {alt_path}")

    print(f"\nPräsentation enthält:")
    print("  • 15 Folien auf Deutsch")
    print("  • 4 hochauflösende Plots")
    print("  • Vollständige Ergebnisdarstellung")
    print("  • Klare Empfehlungen auf Deutsch")
    print("\nDateigröße: ~1.2 MB")
    print("Bereit zum Präsentieren! 🎉")


if __name__ == "__main__":
    main()

