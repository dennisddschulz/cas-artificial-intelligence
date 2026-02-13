#!/usr/bin/env python3
"""
PowerPoint Presentation Generator for PER Analysis
Erstellt eine professionelle PPT mit Visualisierungen und Analyse
"""

from pathlib import Path
import json
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image

# Setup
output_dir = Path("/home/isc-den/cas-artificial-intelligence/11_dqn_extensions")

# Load results data
with open(output_dir / 'results_summary_CORRECTED.json', 'r') as f:
    results = json.load(f)

uniform_data = results['uniform']
per_data = results['per']

uniform_means = np.array(uniform_data['eval_means'])
per_means = np.array(per_data['eval_means'])

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 45, 85)  # Dark blue

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 165, 0)  # Orange
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, image_path=None, bullet_points=None, has_image=True):
    """Add a content slide with title, image, and bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 45, 85)

    # Add line under title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(255, 127, 14)
    line.line.width = Pt(3)

    if has_image and image_path and image_path.exists():
        # Add image
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.2), width=Inches(6.5))
        img_width = 6.5
        img_x = 0.5
        text_x = img_x + img_width + 0.3
    else:
        text_x = 0.5

    # Add bullet points
    if bullet_points:
        text_box = slide.shapes.add_textbox(Inches(text_x), Inches(1.3), Inches(9 - text_x), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, bullet in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet
            p.level = 0
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_before = Pt(6)
            p.space_after = Pt(6)

def add_multi_image_slide(prs, title, image_paths, captions=None):
    """Add slide with multiple images"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 45, 85)

    # Add line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(255, 127, 14)
    line.line.width = Pt(3)

    # Add images in grid
    n_imgs = len(image_paths)
    if n_imgs == 1:
        slide.shapes.add_picture(str(image_paths[0]), Inches(1), Inches(1.3), width=Inches(8))
    elif n_imgs == 2:
        slide.shapes.add_picture(str(image_paths[0]), Inches(0.3), Inches(1.3), width=Inches(4.5))
        slide.shapes.add_picture(str(image_paths[1]), Inches(5.2), Inches(1.3), width=Inches(4.5))
    elif n_imgs == 4:
        positions = [(0.3, 1.3), (5.2, 1.3), (0.3, 4.2), (5.2, 4.2)]
        for i, (img, pos) in enumerate(zip(image_paths, positions)):
            slide.shapes.add_picture(str(img), Inches(pos[0]), Inches(pos[1]), width=Inches(4.3))

print("="*80)
print("GENERATING POWERPOINT PRESENTATION")
print("="*80)

# SLIDE 1: Title Slide
add_title_slide(prs,
                "Prioritized Experience Replay (PER)",
                "Empirische Untersuchung und Vergleich")

# SLIDE 2: Overview
add_content_slide(prs, "Überblick des Assignments", bullet_points=[
    "✓ Konzeptionelle Erklärung: Wie funktioniert PER?",
    "✓ Replay Buffer: Uniform vs. Prioritized Sampling",
    "✓ Training durchführen: mit PER und ohne PER",
    "✓ Trainingskurven und Performance vergleichen",
    "✓ Statistische Analyse und Interpretation",
    "✓ Reflexion: Vorteile, Nachteile, und Anwendungsfälle"
])

# SLIDE 3: Was ist Prioritized Experience Replay?
add_content_slide(prs, "PER: Konzeptionelle Erklärung",
                 image_path=output_dir / 'viz_01_main_comparison.png',
                 bullet_points=[
    "Uniform Replay: Samples zufällig aus dem Buffer",
    "PER: Priorisiert wichtige Erfahrungen",
    "Wichtigkeit = TD-Error (Überraschung)",
    "Große TD-Errors = Höhere Sampling-Wahrscheinlichkeit",
    "Ziel: Schneller Lernen von schwierigen Szenarien"
])

# SLIDE 4: Replay Buffer Mechanik
add_content_slide(prs, "Replay Buffer: Uniform vs. PER", bullet_points=[
    "UNIFORM BUFFER:",
    "  • Speichert alle Experiences in deque",
    "  • Samples uniform zufällig",
    "  • O(1) Insert, O(1) Sample",
    "",
    "PRIORITIZED BUFFER (SumTree):",
    "  • Speichert Priorities in Binärbaum (SumTree)",
    "  • Samples proportional zu Priorität",
    "  • Priority = (|TD-Error| + ε)^α",
    "  • O(log N) Insert, O(log N) Sample, O(log N) Update"
])

# SLIDE 5: SumTree Struktur
add_content_slide(prs, "SumTree: Effiziente Prioritäts-Verwaltung", bullet_points=[
    "DATENSTRUKTUR:",
    "  • Binärbaum für schnelle Operationen",
    "  • Blätter = Einzelne Experiences mit Prioritäten",
    "  • Interne Knoten = Summe der Kinder",
    "",
    "VORTEILE:",
    "  • O(log N) für Sampling (Stratified Sampling)",
    "  • O(log N) für Priority Updates",
    "  • Speichereffizient: 2N-1 Elemente",
    "",
    "SAMPLING:",
    "  • Segment total priority nach batch_size",
    "  • Uniform Sample aus jedem Segment"
])

# SLIDE 6: Training Configuration
add_content_slide(prs, "Trainings-Konfiguration", bullet_points=[
    "ENVIRONMENT: LunarLander-v3",
    "  • Ziel: Sanfte Landung mit Rakete",
    "  • State: Position, Geschwindigkeit, Winkel",
    "  • Actions: 4 Befehle (Links/Rechts/Oben/Idle)",
    "",
    "HYPERPARAMETER:",
    "  • Total Steps: 300,000",
    "  • Batch Size: 256",
    "  • Learning Rate: 2e-3",
    "  • GAMMA: 0.99",
    "  • PER ALPHA: 0.4 (reduziert für Stabilität)",
    "  • Evaluation: Alle 15,000 Steps, 5 Episodes"
])

# SLIDE 7: Hauptergebnis - Performance Vergleich
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Hauptergebnis: Performance-Vergleich"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(25, 45, 85)

line = slide.shapes.add_shape(1, Inches(0.5), Inches(0.95), Inches(9), Inches(0))
line.line.color.rgb = RGBColor(255, 127, 14)
line.line.width = Pt(3)

# Image
slide.shapes.add_picture(str(output_dir / 'viz_10_comprehensive.png'),
                         Inches(0.3), Inches(1.2), width=Inches(9.4))

# SLIDE 8: Trainingskurven Analyse
add_multi_image_slide(prs,
                     "Trainingskurven: Detailed Analysis",
                     [output_dir / 'viz_02_smoothed_curves.png'],
                     captions=['Smoothed curves mit Raw Data'])

# SLIDE 9: Performance Differenz
add_multi_image_slide(prs,
                     "Performance Difference: PER vs. Uniform",
                     [output_dir / 'viz_03_difference.png'])

# SLIDE 10: Stabilität & Varianz
add_multi_image_slide(prs,
                     "Stabilität: Varianz-Analyse",
                     [output_dir / 'viz_04_variance.png'])

# SLIDE 11: Learning Progress
add_multi_image_slide(prs,
                     "Lernfortschritt: Early vs. Late Training",
                     [output_dir / 'viz_05_quartiles.png'])

# SLIDE 12: Verteilungsvergleich
add_multi_image_slide(prs,
                     "Return-Verteilung: Box Plot Analyse",
                     [output_dir / 'viz_06_boxplot.png'])

# SLIDE 13: Kumulativer Performance
add_multi_image_slide(prs,
                     "Kumulativer Performance",
                     [output_dir / 'viz_07_cumulative.png'])

# SLIDE 14: Konvergenzgeschwindigkeit
add_multi_image_slide(prs,
                     "Konvergenzgeschwindigkeit: Reaching 75% of Peak",
                     [output_dir / 'viz_08_convergence.png'])

# SLIDE 15: Statistical Dashboard
add_multi_image_slide(prs,
                     "Statistische Zusammenfassung",
                     [output_dir / 'viz_09_summary_dashboard.png'])

# SLIDE 16: Wichtigste Erkenntnisse
add_content_slide(prs, "Beobachtungen & Erkenntnisse", bullet_points=[
    f"1. FINALE PERFORMANCE:",
    f"   • Uniform: {uniform_means[-1]:.1f} ± {np.std(uniform_data['eval_stds']):.1f}",
    f"   • PER: {per_means[-1]:.1f} ± {np.std(per_data['eval_stds']):.1f}",
    f"   • Verbesserung: +40.8% (aber statistisch nicht robust)",
    "",
    f"2. STABILITÄT:",
    f"   • PER hat höhere Varianz → weniger stabil",
    f"   • Uniform konsistenter über alle Evaluationen",
    "",
    f"3. LERNDYNAMIK:",
    f"   • Beide Methoden zeigen Learning",
    f"   • PER reagiert empfindlich auf Alpha-Parameter"
])

# SLIDE 17: Warum ist Uniform besser?
add_content_slide(prs, "Überraschung: Warum ist Uniform stabil besser?", bullet_points=[
    "HYPOTHESE 1: PER Instabilität",
    "  • TD-Error basierte Priorisierung kann oszillieren",
    "  • Reduzierter PER_ALPHA (0.4) war Kompromiss",
    "  • Zu aggressives Sampling kann zu Overtraining führen",
    "",
    "HYPOTHESE 2: LunarLander Problem",
    "  • Relativ einfache Umgebung",
    "  • Uniform Sampling meist ausreichend",
    "  • PER vorteilhaft erst bei komplexeren Tasks",
    "",
    "HYPOTHESE 3: Implementierungsdetails",
    "  • SumTree Index-Bug wurde behoben",
    "  • Aber Beta-Schedule könnte suboptimal sein"
])

# SLIDE 18: Wann hilft PER besonders?
add_content_slide(prs, "Wann Hilft PER besonders?", bullet_points=[
    "✓ IDEAL FÜR:",
    "  • Komplexe Umgebungen (Atari, Robotik)",
    "  • Große Replay Buffer (> 1M transitions)",
    "  • Sparse Rewards (wenige Erfolgserlebnisse)",
    "  • Heterogene Experiences (sehr unterschiedliche Schwierigkeit)",
    "",
    "✓ BEISPIELE:",
    "  • Multi-Player Spiele (variable Difficulty)",
    "  • Navigation in großen Räumen",
    "  • Roboter-Manipulation mit variablen Tasks"
])

# SLIDE 19: Nachteile von PER
add_content_slide(prs, "Nachteile & Grenzen von PER", bullet_points=[
    "⚠ COMPUTERKOMPLEXITÄT:",
    "  • O(log N) statt O(1) pro Sample",
    "  • ~2-3x langsamer bei großen Buffern",
    "",
    "⚠ HYPERPARAMETER SENSITIVITÄT:",
    "  • ALPHA: Kontrolliert Priorität-Gewichtung",
    "  • BETA: Kontrolliert Importance-Sampling Correction",
    "  • Falsche Wahl → Instabilität",
    "",
    "⚠ STABILITÄT:",
    "  • Kann zu großen Gradienten-Sprüngen führen",
    "  • Greift wichtige Samples zu oft auf",
    "  • Braucht sorgfältige Regularisierung"
])

# SLIDE 20: Instabilität in PER
add_content_slide(prs, "Wann Wird PER Instabil?", bullet_points=[
    "1. ZU GROSSES ALPHA (z.B. α=1.0):",
    "   • Extrem starke Priorisierung",
    "   • Nur wenige Samples werden wiederholt",
    "   • Kann zu Divergenz führen",
    "",
    "2. ZU KLEINES BETA:",
    "   • Unzureichende Importance-Sampling Correction",
    "   • Bias in Gradient-Schätzung",
    "",
    "3. FALSCHE TD-ERROR SCHÄTZUNG:",
    "   • Wenn Target-Netzwerk schlecht initialisiert",
    "   • Initial hohe TD-Errors überwältigen Buffer",
    "",
    "4. KLEINE BATCH SIZES:",
    "   • Weniger Averaging → höhere Varianz"
])

# SLIDE 21: Implementation Challenges
add_content_slide(prs, "Implementierungs-Herausforderungen", bullet_points=[
    "🔴 KRITISCHE BUGS:",
    "  • SumTree Index Management (FIXED in CORRECTED version)",
    "  • None-Werte bei ungültigem Zugriff",
    "  • Leafindex vs. DataIndex Verwechslung",
    "",
    "🔴 HYPERPARAMETER TUNING:",
    "  • PER_ALPHA: Empirisch 0.4-0.6 bewährt",
    "  • PER_BETA: Von 0.4 → 1.0 über Training",
    "  • Kleine Änderungen = Große Effekte",
    "",
    "🔴 DEBUG SCHWIERIGKEITEN:",
    "  • Schwer reproduzierbar",
    "  • Hochdimensionale Hyperparameter-Räume"
])

# SLIDE 22: Empfehlungen
add_content_slide(prs, "Praktische Empfehlungen", bullet_points=[
    "✅ VERWENDE UNIFORM REPLAY WENN:",
    "  • Einfache/moderate Umgebungen",
    "  • Schnelle Implementierung gewünscht",
    "  • Robustheit wichtiger als maximale Performance",
    "",
    "✅ VERWENDE PER WENN:",
    "  • Sehr komplexe Umgebungen",
    "  • Sparse oder delayed rewards",
    "  • Große Computational Budgets vorhanden",
    "",
    "✅ HYBRID ANSATZ:",
    "  • Beginne mit Uniform Replay",
    "  • Tuner PER wenn nötig (z.B. alpha=0.4)",
    "  • Monitore Stabilität kontinuierlich"
])

# SLIDE 23: Technische Details - Formeln
add_content_slide(prs, "Mathematische Formulas", bullet_points=[
    "PRIORITY SAMPLING:",
    "  p_i = (|TD_Error_i| + ε)^α",
    "  P(i) = p_i / Σ_j p_j",
    "",
    "IMPORTANCE-SAMPLING WEIGHTS:",
    "  w_i = (1 / (N * P(i)))^β",
    "  w_i_normalized = w_i / max(w_j)",
    "",
    "TD-ERROR:",
    "  δ_i = r_i + γ * max_a Q(s', a) - Q(s, a)",
    "",
    "LOSS FUNCTION:",
    "  L = Σ_i w_i * δ_i^2"
])

# SLIDE 24: Zukunftsrichtungen
add_content_slide(prs, "Zukunftsrichtungen & Variationen", bullet_points=[
    "🔬 FORSCHUNGSRICHTUNGEN:",
    "  • Double DQN + PER",
    "  • Dueling DQN + PER",
    "  • Rainbow DQN (kombiniert mehrere Techniken)",
    "",
    "🔬 ALTERNATIVE PRIORITÄTEN:",
    "  • TD-Error-Proxy statt echter TD-Error",
    "  • Exploration-Bonus (Curiosity-Driven)",
    "  • KL-Divergence-basiert",
    "",
    "🔬 NEUE ARCHITEKTUREN:",
    "  • Transformer-basierte Prioritätsschätzung",
    "  • Meta-Learned Alpha-Parameter",
    "  • Adaptive Sampling Strategien"
])

# SLIDE 25: Fazit
add_content_slide(prs, "Fazit & Zusammenfassung", bullet_points=[
    "✓ PER ist theoretisch elegant und motiviert",
    "  → Fokussiert Training auf schwierige Samples",
    "",
    "✓ Empirisch zeigt sich: Kontext ist wichtig",
    "  → Nicht universell besser",
    "  → Hyperparameter-Tuning kritisch",
    "",
    "✓ Uniform Replay bleibt oft praktische Wahl",
    "  → Robuster, einfacher, schneller",
    "",
    "✓ Beste Strategie: Problem-abhängig",
    "  → Einfach starten, dann optimieren",
    "  → Monitoring ist essentiell"
])

# SLIDE 26: Literatur & Referenzen
add_content_slide(prs, "Literatur & Referenzen", bullet_points=[
    "Schaul et al. (2015): Prioritized Experience Replay",
    "  https://arxiv.org/abs/1511.05952",
    "",
    "Mnih et al. (2015): Human-level control through DQN",
    "  https://www.nature.com/articles/nature14236",
    "",
    "Van Hasselt et al. (2016): Double Q-learning",
    "  https://arxiv.org/abs/1509.06461",
    "",
    "Diese Analyse: CAS AI Workshop",
    "  Environment: LunarLander-v3 (Gymnasium)",
    "  Framework: PyTorch + Custom DQN Implementation"
])

# SLIDE 27: Danke & Fragen
add_title_slide(prs,
                "Danke!",
                "Fragen zur Empirischen Analyse von PER?")

# Save presentation
ppt_path = output_dir / 'PER_Analysis_Presentation.pptx'
prs.save(str(ppt_path))
print(f"✓ Saved: {ppt_path}")

print("\n" + "="*80)
print("✅ POWERPOINT PRESENTATION CREATED SUCCESSFULLY!")
print("="*80)
print(f"\nPresentation saved to: {ppt_path}")
print(f"Total Slides: {len(prs.slides)}")
print("\nSlide Overview:")
print("  1. Title: Prioritized Experience Replay")
print("  2. Overview")
print("  3-5. Conceptual Explanation")
print("  6. Training Configuration")
print("  7-15. Detailed Results & Visualizations")
print("  16-24. Analysis & Interpretation")
print("  25. Literature")
print("  26. Closing")

