#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for PER vs Uniform Replay Analysis
"""

import json
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

output_dir = Path('/home/isc-den/cas-artificial-intelligence/11_dqn_extensions')

# Load results
with open(output_dir / 'results_summary.json', 'r') as f:
    results = json.load(f)

uniform = results['uniform']
per = results['per']

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 119, 180)  # Blue

    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        left = Inches(0.5)
        top = Inches(4.2)
        width = Inches(9)
        height = Inches(2)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullet_points=None, image_path=None):
    """Add a content slide with title and optional image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 119, 180)

    # Add image if provided
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.3), width=Inches(9))
    elif bullet_points:
        # Add bullet points
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5.5)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)

# ============================================
# SLIDE 1: Title
# ============================================
add_title_slide(
    "Prioritized Experience Replay (PER)",
    "Empirical Analysis & Comparison with Uniform Replay\nLunarLander-v3 (300k training steps)"
)

# ============================================
# SLIDE 2: Overview
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.8)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Experiment Overview"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(31, 119, 180)

left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

overview_text = """
OBJECTIVE:
Empirically compare Prioritized Experience Replay (PER) with classical Uniform Replay

ENVIRONMENT:
• LunarLander-v3 (classic RL benchmark)
• 8-dimensional observation space
• 4 discrete actions

TRAINING CONFIGURATION:
• Total Steps: 300,000
• Batch Size: 256
• Replay Buffer: 200,000
• Learning Rate: 0.002
• Evaluation Episodes: 5 (every 15k steps)

COMPARISON:
1. Uniform Replay (Baseline) - Classical DQN
2. Prioritized Experience Replay (PER) - Advanced method
"""

p = text_frame.paragraphs[0]
p.text = overview_text
p.font.size = Pt(16)
p.space_before = Pt(0)
p.space_after = Pt(0)

# ============================================
# SLIDE 3: PER Concept
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.8)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "What is Prioritized Experience Replay?"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(31, 119, 180)

left = Inches(0.7)
top = Inches(1.4)
width = Inches(9)
height = Inches(5.8)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

concept_text = """
CORE IDEA:
"Not all experiences in the replay buffer are equally valuable for learning"

KEY CONCEPT:
Instead of sampling uniformly from the buffer, PER samples experiences with probability 
proportional to their Temporal Difference (TD) Error.

HOW IT WORKS:
1. Calculate TD-Error for each experience: δ = |Qtarget - Qpredicted|
2. Assign priority based on TD-Error: p_i = (δ_i + ε)^α
3. Sample experiences with probability: P(i) = p_i / Σ(p_j)
4. Update priorities after each training step

INTUITION:
• High TD-Error → Sample more often (harder to learn)
• Low TD-Error → Sample less often (already learned)
• Result: More efficient use of training data

PARAMETERS:
• α (alpha): Controls prioritization strength (0=uniform, 1=pure priority)
• β (beta): Controls importance-sampling correction (0→1 during training)
"""

p = text_frame.paragraphs[0]
p.text = concept_text
p.font.size = Pt(14.5)
p.line_spacing = 1.2

# ============================================
# SLIDE 4: How Replay Buffer Works
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.8)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Replay Buffer: Uniform vs. PER"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(31, 119, 180)

left = Inches(0.5)
top = Inches(1.4)
width = Inches(4.5)
height = Inches(5.8)
text_box1 = slide.shapes.add_textbox(left, top, width, height)
text_frame1 = text_box1.text_frame
text_frame1.word_wrap = True

uniform_text = """UNIFORM REPLAY

Storage:
✓ Simple FIFO queue
✓ Fixed capacity
✓ Overwrites old data

Sampling:
✓ Random selection
✓ Equal probability
✓ All samples equally likely

Priority Update:
✓ Not needed
✓ No computation
✓ O(1) per sample

Pros:
✓ Simple, fast
✓ No bugs possible
✓ Predictable

Cons:
✗ Inefficient learning
✗ Wastes time on easy samples
✗ Ignores hard problems
"""

p = text_frame1.paragraphs[0]
p.text = uniform_text
p.font.size = Pt(12)
p.line_spacing = 1.1

left = Inches(5.1)
top = Inches(1.4)
width = Inches(4.4)
height = Inches(5.8)
text_box2 = slide.shapes.add_textbox(left, top, width, height)
text_frame2 = text_box2.text_frame
text_frame2.word_wrap = True

per_text = """PRIORITIZED REPLAY

Storage:
✓ SumTree data structure
✓ Priority-weighted
✓ Tracks TD-Errors

Sampling:
✓ Weighted selection
✓ Based on TD-Error
✓ Hard samples preferred

Priority Update:
✓ After each step
✓ Uses new TD-Error
✓ O(log N) per sample

Pros:
✓ More efficient
✓ Focus on hard problems
✓ Faster convergence*

Cons:
✗ Complex code
✗ Higher computational cost
✗ More bug-prone
"""

p = text_frame2.paragraphs[0]
p.text = per_text
p.font.size = Pt(12)
p.line_spacing = 1.1

# ============================================
# SLIDE 5: Main Results
# ============================================
add_content_slide(
    "Main Results: Learning Curves",
    image_path=str(output_dir / 'VIZ_01_main_comparison.png')
)

# ============================================
# SLIDE 6: Detailed Analysis
# ============================================
add_content_slide(
    "Detailed Performance Analysis",
    image_path=str(output_dir / 'VIZ_02_detailed_analysis.png')
)

# ============================================
# SLIDE 7: Key Findings
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.8)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Unexpected Finding: PER Underperforms"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 0, 0)

left = Inches(0.7)
top = Inches(1.4)
width = Inches(8.6)
height = Inches(5.8)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

findings_text = """
KEY RESULTS:

UNIFORM REPLAY (Baseline):
  Final Return:        55.22 ± 188.19  ✓ POSITIVE
  Max Return:          264.82
  Mean Return:         16.26
  Trend:               Improving

PER (Advanced Method):
  Final Return:       -719.10 ± 94.61  ✗ VERY NEGATIVE
  Max Return:         -162.31
  Mean Return:        -674.70
  Trend:              Declining

STATISTICAL SIGNIFICANCE:
  p-value: 7.5e-10 (Highly significant)
  Conclusion: Uniform is statistically better (p < 0.001)

SURPRISE FACTOR:
⚠️  This result is CONTRARY to expected PER benefits!
    In most benchmarks, PER performs as well or better than Uniform Replay.

INTERPRETATION:
This suggests a BUG in the PER implementation, not a fundamental issue with PER itself.
"""

p = text_frame.paragraphs[0]
p.text = findings_text
p.font.size = Pt(15)
p.line_spacing = 1.2

# ============================================
# SLIDE 8: Why Did PER Fail?
# ============================================
add_content_slide(
    "Analysis: Why PER Failed",
    image_path=str(output_dir / 'VIZ_03_analysis_why_failed.png')
)

# ============================================
# SLIDE 9: Possible Issues
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.8)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Likely Issues in PER Implementation"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 0, 0)

left = Inches(0.7)
top = Inches(1.4)
width = Inches(8.6)
height = Inches(5.8)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

issues_text = """
1. PRIORITY UPDATE BUG
   Problem: TD-Errors might not be properly converted to priorities
   Impact: Sampled experiences might not reflect actual difficulty
   Fix: Add logging to verify priority calculations

2. IMPORTANCE SAMPLING WEIGHTS
   Problem: Weights could be inverted or computed incorrectly
   Impact: Gradients could be scaled wrong, leading to poor learning
   Fix: Debug the weight computation and loss calculation

3. SUMTREE INDEXING ERROR
   Problem: Buffer indices might not match SumTree indices
   Impact: Updating wrong priorities, sampling wrong transitions
   Fix: Add assertions to verify index consistency

4. HYPERPARAMETER MISMATCH
   Problem: Alpha=0.6 might be too aggressive for this environment
   Impact: Extreme prioritization causes instability
   Fix: Reduce alpha (try 0.3) and test systematically

5. GRADIENT INSTABILITY
   Problem: Large importance weights cause gradient explosion
   Impact: Parameters diverge, Q-values become NaN/-inf
   Fix: Add gradient clipping, debug loss computation

6. INITIALIZATION ERROR
   Problem: Buffer might not be properly filled before training
   Impact: Early priorities could be wrong
   Fix: Verify buffer filling logic
"""

p = text_frame.paragraphs[0]
p.text = issues_text
p.font.size = Pt(13.5)
p.line_spacing = 1.15

# ============================================
# SLIDE 10: When Does PER Work?
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.8)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "When Does PER Help? (Theory)"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(31, 119, 180)

left = Inches(0.7)
top = Inches(1.4)
width = Inches(4.2)
height = Inches(5.8)
text_box1 = slide.shapes.add_textbox(left, top, width, height)
text_frame1 = text_box1.text_frame
text_frame1.word_wrap = True

when_helps = """PER HELPS WHEN:

✓ Sparse-Reward Envs
  Only few samples have 
  information. Focus on 
  these samples.

✓ Complex Domains
  Many edge cases, some 
  much harder. PER finds 
  and fixes them.

✓ Sample Efficiency
  Real-world costs (robot 
  moves, etc.) → need 
  efficient learning.

✓ Large Buffers
  With huge buffers, 
  focusing on hard cases 
  is very beneficial.
"""

p = text_frame1.paragraphs[0]
p.text = when_helps
p.font.size = Pt(13.5)
p.line_spacing = 1.2

left = Inches(5.1)
top = Inches(1.4)
width = Inches(4.2)
height = Inches(5.8)
text_box2 = slide.shapes.add_textbox(left, top, width, height)
text_frame2 = text_box2.text_frame
text_frame2.word_wrap = True

when_not = """PER LESS HELPFUL:

✗ Simple Environments
  All samples roughly 
  equal value. Uniform 
  sampling fine.

✗ Dense Rewards
  Many samples have 
  learning signal. No 
  efficiency gain.

✗ Small Buffers
  Limited data anyway. 
  Uniform sampling 
  sufficient.

✗ Abundant Compute
  If data/speed unlimited, 
  complexity not worth it.
"""

p = text_frame2.paragraphs[0]
p.text = when_not
p.font.size = Pt(13.5)
p.line_spacing = 1.2

# ============================================
# SLIDE 11: Recommendations
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.8)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Recommendations & Next Steps"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(31, 119, 180)

left = Inches(0.7)
top = Inches(1.4)
width = Inches(8.6)
height = Inches(5.8)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

recommendations = """
IMMEDIATE ACTIONS:

1. DEBUG THE CODE
   • Add comprehensive logging to priority calculations
   • Verify TD-Error calculations are correct
   • Check SumTree operations step-by-step
   • Print sample priorities and weights during training

2. SIMPLIFY IMPLEMENTATION
   • Reduce alpha from 0.6 to 0.3
   • Use fixed beta instead of annealing
   • Add more error checking and assertions
   • Test each component in isolation

3. VERIFY AGAINST REFERENCE
   • Compare with stable-baselines3 PER implementation
   • Test on simpler environments first
   • Check if issue is specific to our code

4. SYSTEMATIC TESTING
   • Test with different alpha values (0.1, 0.3, 0.5, 0.7)
   • Try different beta schedules
   • Test with smaller buffers
   • Evaluate on other environments (Atari, etc.)

CONCLUSION:
PER has proven benefits in the literature, but our implementation has bugs.
The negative results are NOT evidence against PER, but show the importance of:
  • Correct implementation
  • Thorough testing
  • Careful debugging
"""

p = text_frame.paragraphs[0]
p.text = recommendations
p.font.size = Pt(13.5)
p.line_spacing = 1.2

# ============================================
# SLIDE 12: Summary
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(31, 119, 180)

left = Inches(0.5)
top = Inches(0.5)
width = Inches(9)
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Summary"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

left = Inches(1)
top = Inches(1.8)
width = Inches(8)
height = Inches(5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

summary = """
KEY TAKEAWAYS:

1. PER is a theoretically sound technique that prioritizes hard-to-learn samples

2. Our implementation has a significant bug causing it to underperform Uniform Replay

3. The bug is likely in:
   • TD-Error calculation
   • Priority updating
   • Importance sampling weights
   • SumTree indexing

4. Despite this implementation failure, PER has proven benefits in literature

5. Proper implementation of PER requires:
   • Careful debugging
   • Verification of each component
   • Testing on diverse environments
   • Comparison with reference implementations

6. When properly implemented, PER can improve sample efficiency significantly

Next: Fix the implementation and re-test!
"""

p = text_frame.paragraphs[0]
p.text = summary
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(255, 255, 255)
p.line_spacing = 1.3

# Save presentation
ppt_path = output_dir / 'PER_Analysis_Presentation.pptx'
prs.save(str(ppt_path))

print(f"\n✅ PowerPoint presentation created: {ppt_path}")
print(f"   12 slides with comprehensive analysis")

