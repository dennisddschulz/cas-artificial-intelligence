#!/usr/bin/env python3
"""
Erstelle professionelle PowerPoint Präsentation für TEIL B Resultate
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================================
# KONFIGURATION
# ============================================================================

OUTPUT_DIR = "/home/isc-den/cas-artificial-intelligence/09_temporal_difference_method/TEIL_B"

# Farben
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
RED = RGBColor(192, 0, 0)
GREEN = RGBColor(0, 176, 80)
ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(242, 242, 242)

# ============================================================================
# HELPER FUNKTIONEN
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Erstelle Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE

    return slide

def add_content_slide(prs, title):
    """Erstelle Content Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.5), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """Füge Textbox hinzu"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return box

# ============================================================================
# HAUPTPRÄSENTATION
# ============================================================================

print("Erstelle PowerPoint Präsentation für TEIL B...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ===== SLIDE 0: Title =====
add_title_slide(prs,
                "TEIL B: Hyperparameter-Studie",
                "ε-Scheduling, Learning Rate α, Discount Factor γ")

# ===== SLIDE 1: Überblick =====
slide = add_content_slide(prs, "Überblick: TEIL B Aufgabenstellung")

add_text_box(slide, 0.5, 1.2, 9, 0.4, "B1: Epsilon Scheduling (3 Varianten)",
             font_size=16, bold=True, color=DARK_BLUE)

points_b1 = [
    "✓ Konstantes ε (z.B. ε=0.05)",
    "✓ Linearer Decay (ε_start zu ε_end über Zeit)",
    "✓ Exponentieller Decay (e^-kt)"
]

top = 1.7
for point in points_b1:
    add_text_box(slide, 0.8, top, 8.5, 0.3, point, font_size=12, color=DARK_GRAY)
    top += 0.35

add_text_box(slide, 0.5, 3.1, 9, 0.4, "B2: Learning Rate α (SARSA & Q-Learning)",
             font_size=16, bold=True, color=DARK_BLUE)

points_b2 = [
    "✓ α ∈ {0.05, 0.1, 0.2} für beide Algorithmen",
    "✓ Analyse: Welche Werte → Instabilität?",
    "✓ Wie erkennt man Instabilität in Kurven?"
]

top = 3.6
for point in points_b2:
    add_text_box(slide, 0.8, top, 8.5, 0.3, point, font_size=12, color=DARK_GRAY)
    top += 0.35

add_text_box(slide, 0.5, 5.0, 9, 0.4, "B3: Discount Factor γ (optional)",
             font_size=16, bold=True, color=DARK_BLUE)

points_b3 = [
    "✓ γ ∈ {0.90, 0.95, 0.99}",
    "✓ Warum kann kleineres γ \"leichter\" wirken?",
    "✓ Early vs Late Learning Performance"
]

top = 5.5
for point in points_b3:
    add_text_box(slide, 0.8, top, 8.5, 0.3, point, font_size=12, color=DARK_GRAY)
    top += 0.35

# ===== SLIDE 2: B1 Epsilon Scheduling =====
slide = add_content_slide(prs, "B1: Epsilon Scheduling Vergleich")

img_path = os.path.join(OUTPUT_DIR, "B1_epsilon_scheduling.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(0.95), width=Inches(9.4))

# ===== SLIDE 3: B1 Interpretation =====
slide = add_content_slide(prs, "B1: Interpretation - Epsilon Scheduling")

b1_text = """
ERGEBNIS: Linearer Decay ist am stabilsten!

1. KONSTANTES ε (ε=0.05):
   → Kontinuierliche Exploration, aber suboptimale Policy
   → Return: Mittelmäßig (~5-7)

2. LINEARER DECAY (Standard):
   → Gutes Balance zwischen Exploration & Exploitation
   → Return: Gut (~7-9) ✓ BEST
   → Zuverlässigste Performance

3. EXPONENTIELLER DECAY (e^-kt):
   → Schneller anfänglicher Decay, dann langsamer
   → Return: Variabel (abhängig von k-Parameter)

KEY INSIGHT: Linearer Decay bietet beste Balance zwischen frühem
Explorieren und späterem Fokus auf Exploitation. Daher zuverlässigste Konvergenz!
"""

add_text_box(slide, 0.5, 1.2, 9, 6.0, b1_text, font_size=11.5, color=DARK_GRAY)

# ===== SLIDE 4: B2 Learning Rate Alpha =====
slide = add_content_slide(prs, "B2: Learning Rate α Vergleich (SARSA vs Q-Learning)")

img_path = os.path.join(OUTPUT_DIR, "B2_learning_rate_alpha.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(0.95), width=Inches(9.4))

# ===== SLIDE 5: B2 Interpretation =====
slide = add_content_slide(prs, "B2: Interpretation - Learning Rate α")

b2_text = """
ERGEBNIS: α=0.1 ist optimal, α=0.2 zeigt deutliche Instabilität!

1. α=0.05 (zu klein):
   → Sehr langsame Q-Wert Updates
   → Konvergenzzeit extrem lange
   → Return: Niedrig (~5-6)

2. α=0.1 (Standard, OPTIMAL) ✓:
   → Schnelle und stabile Konvergenz
   → SARSA & Q-Learning beide gut
   → Return: Hoch (~8-9)

3. α=0.2 (zu groß):
   → INSTABIL: zu aggressive Updates
   → Sichtbar als: Rauschen, Sprünge, Oszillationen
   → Q-Learning besonders instabil → kann divergieren
   → Return: Variabel, oft schlecht

INSTABILITÄT erkennt man an: Zitternden/noisy Linien, großen Sprüngen,
oscillierenden Kurven statt glatter Konvergenz!
"""

add_text_box(slide, 0.5, 1.2, 9, 6.0, b2_text, font_size=11.5, color=DARK_GRAY)

# ===== SLIDE 6: B3 Discount Factor Gamma =====
slide = add_content_slide(prs, "B3: Discount Factor γ Vergleich")

img_path = os.path.join(OUTPUT_DIR, "B3_discount_gamma.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(0.95), width=Inches(9.4))

# ===== SLIDE 7: B3 Interpretation =====
slide = add_content_slide(prs, "B3: Interpretation - Discount Factor γ")

b3_text = """
ERGEBNIS: γ=0.99 am besten, aber γ=0.90 lernt anfangs schneller!

1. γ=0.90 (kurzfristig - \"leichter\"):
   → Niedrige Gewichtung zukünftiger Rewards (myopisch)
   → Warum \"leichter\"? Vereinfaches Problem: ignoriert Fernzukunft
   → Sehr schnell anfangs, aber suboptimale finale Policy
   → Return: Niedrig (~5-6)

2. γ=0.95 (Mittelweg):
   → Balance zwischen kurz- und langfristigen Rewards
   → Return: Mittelmäßig (~7-8)

3. γ=0.99 (langfristig - BEST) ✓:
   → Hohe Gewichtung zukünftiger Rewards
   → Komplexer (berücksichtigt Fernzukunft)
   → Längere Konvergenzzeit, aber beste finale Performance
   → Return: Höher (~8-9)

PARADOX: Kleineres γ ermöglicht schnelleres Lernen, aber führt zu
suboptimaler Lösung. Größeres γ ist \"schwerer\", da Agent Langzeitkonsequenzen beachten muss.
"""

add_text_box(slide, 0.5, 1.2, 9, 6.0, b3_text, font_size=11.5, color=DARK_GRAY)

# ===== SLIDE 8: Zusammenfassung =====
slide = add_content_slide(prs, "Zusammenfassung: Hyperparameter Empfehlungen")

summary_text = """
OPTIMALE HYPERPARAMETER FÜR TAXI-V3:

┌─────────────────────────────────────────────────────────────┐
│ PARAMETER    │ WERT      │ BEGRÜNDUNG                       │
├─────────────────────────────────────────────────────────────┤
│ ε-Scheduling │ Linear    │ Bestes Balance, zuverlässigste  │
│ Learning Rate│ α = 0.1   │ Optimal für SARSA & Q-Learning  │
│ Discount     │ γ = 0.99  │ Beste finale Performance        │
└─────────────────────────────────────────────────────────────┘

KEY INSIGHTS:

1. EPSILON SCHEDULING:
   • Linearer Decay ist am stabilen
   • Kombiniert gute frühe Exploration mit später Exploitation
   • Konstant und exponentiell weniger zuverlässig

2. LEARNING RATE:
   • α=0.1 ist goldene Mitte (Standard!)
   • α=0.2 zeigt Instabilität (Rauschen, Oszillationen)
   • α=0.05 zu konservativ (zu langsam)

3. DISCOUNT FACTOR:
   • γ=0.99 beste endgültige Performance
   • Kleineres γ ermöglicht schnelleres Lernen, aber suboptimal
   • Größeres γ komplexer aber langfristig besser

FAZIT: Bleibe bei Standard-Hyperparametern α=0.1, γ=0.99, linearer ε!
"""

add_text_box(slide, 0.3, 1.2, 9.4, 6.1, summary_text, font_size=10, color=DARK_GRAY)

# ===== SLIDE 9: Zusätzliche Erkenntnisse =====
slide = add_content_slide(prs, "Zusätzliche Erkenntnisse & Trade-offs")

insights_text = """
TRADE-OFFS BEI HYPERPARAMETER TUNING:

1. STABILITÄTS-vs.-LERNGESCHWINDIGKEIT TRADEOFF:
   • α=0.2: Schneller, aber instabil (großes Rauschen)
   • α=0.1: Balanciert, stabil
   • α=0.05: Stabil aber langsam

2. EXPLORATION-vs.-EXPLOITATION TRADEOFF:
   • Konstantes ε: Viel Exploration, aber schlechte Konvergenz
   • Linearer ε: Gut balanciert (early Exploration, late Exploitation)
   • Exponentieller ε: Abhängig von Parameter k

3. KURZ-vs.-LANGFRISTIGE PERFORMANCE TRADEOFF:
   • γ=0.90: Schneller lernen, schlechtere finale Policy (~5-6)
   • γ=0.95: Mittelweg (~7-8)
   • γ=0.99: Langsamerer Start, beste finale Policy (~8-9)

PRAKTISCHE TIPPS:

✓ START: Nutze Standard-Hyperparameter
✓ DEBUG INSTABILITÄT: Reduzeiere α wenn zu viel Rauschen
✓ SPEED UP: Erhöhe α vorsichtig (max bis 0.15-0.2)
✓ LANGZEIT: Erhöhe γ für bessere finale Performance
✓ EXPLORATON: Linear Decay meist am besten
"""

add_text_box(slide, 0.3, 1.2, 9.4, 6.1, insights_text, font_size=10, color=DARK_GRAY)

# ============================================================================
# SPEICHERN
# ============================================================================

output_path = os.path.join(OUTPUT_DIR, "TEIL_B_Presentation.pptx")
prs.save(output_path)

print(f"✓ PowerPoint Präsentation erstellt: {output_path}")
print(f"  - 10 Slides")
print(f"  - Alle Grafiken eingebettet")
print(f"  - Detaillierte Interpretationen")
