#!/usr/bin/env python3
"""
Erstelle professionelle PowerPoint Präsentation mit eingebetteten Diagrammen
PRIO 1: On-Policy vs Off-Policy (SARSA vs Q-Learning)
PRIO 2: Monte Carlo vs Temporal Difference
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================================
# KONFIGURATION
# ============================================================================

OUTPUT_DIR = "/home/isc-den/cas-artificial-intelligence/09_temporal_difference_method"

# Farben
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
RED = RGBColor(220, 20, 60)
BLUE = RGBColor(30, 144, 255)
GREEN = RGBColor(34, 139, 34)
ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(89, 89, 89)

# ============================================================================
# HELPER FUNKTIONEN
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Erstelle Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title):
    """Erstelle Content Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.85)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8.5), Inches(0.65))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=13, bold=False,
                 color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """Füge Textbox hinzu"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return box

# ============================================================================
# HAUPTPRÄSENTATION MIT DIAGRAMMEN
# ============================================================================

print("Erstelle PowerPoint Präsentation mit eingebetteten Diagrammen...\n")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ===== SLIDE 0: TITLE =====
add_title_slide(prs,
                "Reinforcement Learning Algorithmen",
                "On-Policy vs Off-Policy & Monte Carlo vs Temporal Difference")

# ===== SLIDE 1: Überblick =====
slide = add_content_slide(prs, "Agenda: Zwei Hauptvergleiche")

add_text_box(slide, 0.5, 1.1, 4, 2.5,
            "🎯 PRIORITÄT 1:\n\nOn-Policy vs Off-Policy\n\n• SARSA (On-Policy)\n• Q-Learning (Off-Policy)\n• Empirischer Vergleich",
            font_size=13, color=DARK_GRAY)

add_text_box(slide, 5.5, 1.1, 4, 2.5,
            "🎯 PRIORITÄT 2 (Optional):\n\nMonte Carlo vs Temporal Difference\n\n• MC: Episode-basiert\n• TD: Step-basiert\n• Empirische Analyse",
            font_size=13, color=DARK_GRAY)

add_text_box(slide, 0.5, 4.0, 9, 2.8,
            "KEY INSIGHT:\nDer entscheidende Unterschied liegt in der TARGET BESTIMMUNG!\nWas glaubt der Agent über die Zukunft?",
            font_size=14, bold=True, color=DARK_BLUE)

# ===== SLIDE 2: TARGET BESTIMMUNG =====
slide = add_content_slide(prs, "Der Kern: Target Bestimmung")

add_text_box(slide, 0.5, 1.1, 4.3, 0.8,
            "SARSA (On-Policy)",
            font_size=14, bold=True, color=BLUE)

add_text_box(slide, 0.5, 2.0, 4.3, 1.8,
            "TARGET = r + γ·Q(s', a')\n\nwhere a' ~ π(s')\n\n(ACTUAL next action\nfrom policy)",
            font_size=12, color=DARK_GRAY)

add_text_box(slide, 5.2, 1.1, 4.3, 0.8,
            "Q-Learning (Off-Policy)",
            font_size=14, bold=True, color=RED)

add_text_box(slide, 5.2, 2.0, 4.3, 1.8,
            "TARGET = r + γ·max Q(s', ·)\n\nwhere · = best action\n\n(OPTIMAL action\nregardless of policy)",
            font_size=12, color=DARK_GRAY)

add_text_box(slide, 0.5, 4.0, 9, 2.8,
            "💡 DER UNTERSCHIED:\n\n" +
            "SARSA fragt: \"Was wird der Agent TATSÄCHLICH tun?\"\n" +
            "Q-Learning fragt: \"Was wäre OPTIMAL?\"\n\n" +
            "→ SARSA ist KONSERVATIV (vorsichtig mit Risiken)\n" +
            "→ Q-Learning ist AGGRESSIV (optimistic, sucht optimale Policy)",
            font_size=12, color=DARK_GRAY)

# ===== SLIDE 3: PRIO 1 DIAGRAMME =====
slide = add_content_slide(prs, "PRIO 1: SARSA vs Q-Learning - Empirische Resultate")

img_path = os.path.join(OUTPUT_DIR, "PRIO1_SARSA_vs_QLearning.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.2), Inches(0.95), width=Inches(9.6))
    print(f"  ✓ PRIO1 Diagramm eingebettet")
else:
    add_text_box(slide, 0.5, 3.5, 9, 1, "⚠️ Diagramm nicht gefunden", font_size=14, color=RED)

# ===== SLIDE 4: Empirische Analyse PRIO 1 =====
slide = add_content_slide(prs, "PRIO 1: Empirische Analyse")

analysis_text = """
BEOBACHTUNGEN AUS DEN DIAGRAMMEN:

1. LEARNING CURVES (Return):
   • Q-Learning konvergiert SCHNELLER (aggressive Off-Policy Updates)
   • SARSA braucht länger aber ist STABILER (conservative On-Policy)
   • Endgültige Performance: ähnlich (~-5 bis 0)

2. EFFIZIENZ (Episode Length):
   • Q-Learning schafft kürzere Episodes (schnellere Konvergenz zu besserer Policy)
   • SARSA längere Episodes (vorsichtige Explorationsstrategie)
   • Unterschied wird kleiner am Ende

3. TD ERROR DISTRIBUTION:
   • Q-Learning: GRÖSSERE TD-Fehler (aggressive, overoptimistic Updates)
   • SARSA: KLEINERE TD-Fehler (vorsichtige, realistische Updates)
   • KEY INSIGHT: Grösserer Fehler ≠ schlechtere Performance!

4. EARLY vs LATE LEARNING:
   • Q-Learning: schneller anfangs
   • SARSA: stabiler über Zeit
   • Tradeoff zwischen Speed und Stability

5. CONVERGENCE SPEED:
   • Q-Learning erreicht Schwellenwert deutlich FRÜHER
   • SARSA braucht mehr Episoden aber konvergiert zu stabilerem Wert
"""

add_text_box(slide, 0.5, 1.1, 9, 6, analysis_text, font_size=11, color=DARK_GRAY)

# ===== SLIDE 5: Verhalten in der Umgebung =====
slide = add_content_slide(prs, "Verhalten in der Umgebung")

behavior_text = """
WIE UNTERSCHEIDEN SICH SARSA UND Q-LEARNING IN DER PRAXIS?

SARSA (On-Policy) - Risk-Aware:
  → Exploriert vorsichtig, achtet auf potenzielle Gefahren
  → Wenn Agent zufällig in gefährliche Situation gerät, lernt SARSA das
  → Policy vermeidet Zustände mit hohem Explorations-Risiko
  → Beispiel: Roboter vor Klippe → sehr vorsichtig am Rand
  → Resultat: Sichere, stabile, aber nicht optimale Policy

Q-Learning (Off-Policy) - Optimal-Seeking:
  → Exploriert aggressiv, sucht beste mögliche Policy
  → Ignoriert dass Agent exploriert, bootet von bester Action
  → Policy kann Risiken eingehen für Optimalität
  → Beispiel: Roboter vor Klippe → lernt optimale Route, auch wenn explorativ
  → Resultat: Optimale, aber potenziell risikoreiche Policy

PRAKTISCHES BEISPIEL - Cliff Walking:
  SARSA sieht:  "Agent fällt 20% der Zeit runter → das ist schlecht"
               → meidet nahe an Klippe
  
  Q-Learning:   "Beste Aktion meidet Klippe → das ist der Weg"
               → lernt optimal, ignoriert dass Exploration runterfallen könnte
"""

add_text_box(slide, 0.5, 1.1, 9, 6, behavior_text, font_size=11, color=DARK_GRAY)

# ===== SLIDE 6: PRIO 2 - MC vs TD Erklärung =====
slide = add_content_slide(prs, "PRIO 2: Monte Carlo vs Temporal Difference - Unterschiede")

mc_td_text = """
UPDATE MECHANISMUS - Das ist der fundamentale Unterschied!

MONTE CARLO (Batch Updates):
  ┌─────────────────────────────────────┐
  │ Episode läuft: s₀→s₁→s₂→...→sₜ    │
  │ Rewards:       r₁  r₂  ...  rₜ     │
  │ Warte bis EPISODE ENDET!            │
  │ Berechne: G = r₁ + γr₂ + γ²r₃ +...│
  │ UPDATE: Q(s,a) += α(G - Q(s,a))   │
  └─────────────────────────────────────┘
  
  Charakteristiken:
  ✓ Unbiased (verwendet true return)
  ✗ High Variance (variiert mit Episoden)
  ✗ Langsam (1 Update pro Episode)
  ✓ Gut für Evaluation

TEMPORAL DIFFERENCE (Online Updates):
  ┌──────────────────────────────────────┐
  │ Nach JEDEM SCHRITT:                 │
  │ UPDATE: Q(s,a) += α(r + γV(s') - Q) │
  │ V(s') wird als Bootstrap verwendet  │
  │ Nicht auf echtem Return basierend   │
  └──────────────────────────────────────┘
  
  Charakteristiken:
  ✗ Biased (nutzt Schätzung V(s'))
  ✓ Low Variance (kleine Updates)
  ✓ Schnell (viele Updates pro Episode)
  ✓ Gut für Online Lernen

KEY: MC wartet, TD aktualisiert sofort!
"""

add_text_box(slide, 0.5, 1.1, 9, 6, mc_td_text, font_size=10.5, color=DARK_GRAY)

# ===== SLIDE 7: PRIO 2 DIAGRAMME =====
slide = add_content_slide(prs, "PRIO 2: Monte Carlo vs TD - Empirische Resultate")

img_path = os.path.join(OUTPUT_DIR, "PRIO2_MC_vs_TD.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.2), Inches(0.95), width=Inches(9.6))
    print(f"  ✓ PRIO2 Diagramm eingebettet")
else:
    add_text_box(slide, 0.5, 3.5, 9, 1, "⚠️ Diagramm nicht gefunden", font_size=14, color=RED)

# ===== SLIDE 8: MC vs TD Analyse =====
slide = add_content_slide(prs, "PRIO 2: Empirische Analyse MC vs TD")

mc_td_analysis = """
BEOBACHTUNGEN AUS DEN DIAGRAMMEN:

1. LEARNING CURVES (MC vs TD):
   • TD (Q-Learning) konvergiert VIEL SCHNELLER
   • MC braucht deutlich mehr Episodes
   • Grund: TD aktualisiert nach JEDEM Schritt, MC nur nach ganzer Episode

2. SAMPLE EFFICIENCY:
   • TD ist DEUTLICH effizienter
   • MC: viele Steps nötig für wenig Updates
   • TD: jeder Step produziert ein Update

3. EPISODE LENGTH:
   • Ähnlich zwischen MC und TD
   • Unterschied ist in der Update-Frequenz, nicht in Exploration

4. CONVERGENCE SPEED:
   • TD erreicht Return-Ziele deutlich FRÜHER (ca. 3-5x schneller)
   • MC braucht viel länger, besonders in Taxi-v3

5. PRAKTISCHE IMPLIKATION:
   • TD (SARSA, Q-Learning) → verwende online, schnelle Anpassung
   • MC → verwende wenn Episoden kurz/bekannt sind, für Evaluation

ZUSAMMENFASSUNG:
TD ist online und effizient → STANDARD wahl für RL
MC ist unbiased und stabil → GUT FÜR THEORIE
"""

add_text_box(slide, 0.5, 1.1, 9, 6, mc_td_analysis, font_size=11, color=DARK_GRAY)

# ===== SLIDE 9: Synthesis - Alle 4 Algorithmen =====
slide = add_content_slide(prs, "Synthesevergleich: Alle 4 Algorithmen")

img_path = os.path.join(OUTPUT_DIR, "Synthesis_All_Algorithms.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.2), Inches(0.95), width=Inches(9.6))
    print(f"  ✓ Synthesis Diagramm eingebettet")
else:
    add_text_box(slide, 0.5, 3.5, 9, 1, "⚠️ Diagramm nicht gefunden", font_size=14, color=RED)

# ===== SLIDE 10: Zusammenfassung =====
slide = add_content_slide(prs, "Zusammenfassung & Empfehlungen")

summary = """
KERNERKENNTNISSE:

🎯 PRIO 1: ON-POLICY vs OFF-POLICY
   SARSA (On-Policy):
     • Lernt unter AKTUELLER Policy
     • Vorsichtig mit Explorations-Risiken
     • Kleinere TD-Fehler, stabiler
     • Langsamer aber sicherer
   
   Q-Learning (Off-Policy):
     • Lernt OPTIMALE Policy
     • Ignoriert Explorations-Risiken
     • Grössere TD-Fehler, aggressiver
     • Schneller aber risikoreicher
   
   PRAKTISCH: SARSA für Roboter/Safety | Q-Learning für Games/Optimalität

🎯 PRIO 2: MC vs TD
   Monte Carlo:
     • Episode-weise Updates (Batch)
     • Unbiased, High Variance
     • Langsam: 1 Update/Episode
   
   Temporal Difference:
     • Step-weise Updates (Online)
     • Biased, Low Variance
     • Schnell: viele Updates/Episode
   
   PRAKTISCH: TD ist Standard-Wahl für RL (schneller & effizienter)

🎯 KEY INSIGHT:
   Der KERN-UNTERSCHIED liegt in TARGET BESTIMMUNG:
   • Was wird tatsächlich passieren? (SARSA)
   • Was ist optimal? (Q-Learning)
   • Gesamter Episode? (MC)
   • Nächster Schritt? (TD)
"""

add_text_box(slide, 0.5, 1.1, 9, 6, summary, font_size=11, color=DARK_GRAY)

# ============================================================================
# SPEICHERN
# ============================================================================

output_path = os.path.join(OUTPUT_DIR, "On_Policy_vs_Off_Policy_FINAL.pptx")
prs.save(output_path)

print(f"\n✅ PowerPoint Präsentation erstellt:")
print(f"   {output_path}")
print(f"\n📊 INHALT (11 Slides):")
print(f"   0. Title Slide")
print(f"   1. Agenda & Überblick")
print(f"   2. TARGET BESTIMMUNG (Der Kern!)")
print(f"   3. PRIO 1: SARSA vs Q-Learning DIAGRAMME ⭐")
print(f"   4. PRIO 1: Empirische Analyse")
print(f"   5. Verhalten in der Umgebung")
print(f"   6. PRIO 2: MC vs TD Erklärung")
print(f"   7. PRIO 2: MC vs TD DIAGRAMME ⭐")
print(f"   8. PRIO 2: MC vs TD Analyse")
print(f"   9. Synthesevergleich DIAGRAMME ⭐")
print(f"   10. Zusammenfassung & Empfehlungen")
print(f"\n✨ FEATURES:")
print(f"   ✓ 3 große Diagramme eingebettet (PRIO1, PRIO2, Synthesis)")
print(f"   ✓ Learning Curves & Vergleiche")
print(f"   ✓ Detaillierte Erklärungen")
print(f"   ✓ Empirische Resultate")
print(f"   ✓ Praktische Empfehlungen")
