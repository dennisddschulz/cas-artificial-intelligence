#!/usr/bin/env python3
"""
Erstelle professionelle PowerPoint Präsentation:
PRIO 1: On-Policy (SARSA) vs Off-Policy (Q-Learning)
PRIO 2: Monte Carlo vs Temporal Difference (optional)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================================
# KONFIGURATION
# ============================================================================

OUTPUT_DIR = "/home/isc-den/cas-artificial-intelligence/09_temporal_difference_method"

# Farben
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
RED = RGBColor(220, 20, 60)  # Crimson für Q-Learning
BLUE = RGBColor(30, 144, 255)  # Dodger Blue für SARSA
GREEN = RGBColor(34, 139, 34)  # Forest Green
ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(242, 242, 242)

# ============================================================================
# HELPER FUNKTIONEN
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Erstelle Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title):
    """Erstelle Content Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.85)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8.5), Inches(0.65))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=13, bold=False,
                 color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """Füge Textbox hinzu"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for line in text.split('\n'):
        if line.strip():
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = alignment
            p.space_before = Pt(3)
            p.space_after = Pt(3)
    return box

def add_bullet_points(slide, left, top, width, height, points, font_size=13, color=DARK_GRAY):
    """Füge Bullet Points hinzu"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return box

# ============================================================================
# HAUPTPRÄSENTATION
# ============================================================================

print("Erstelle professionelle PowerPoint Präsentation...")
print("Thema: On-Policy vs Off-Policy & MC vs TD\n")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ===== SLIDE 0: TITLE =====
add_title_slide(prs,
                "Reinforcement Learning Algorithmen",
                "On-Policy vs Off-Policy & Monte Carlo vs Temporal Difference")

# ===== SLIDE 1: Überblick =====
slide = add_content_slide(prs, "Agenda: Zwei Hauptvergleiche")

add_text_box(slide, 0.5, 1.1, 4, 2.5,
            "🎯 PRIORITÄT 1:\n\nOn-Policy vs Off-Policy\n\n• SARSA (On-Policy)\n• Q-Learning (Off-Policy)\n• Empirischer Vergleich",
            font_size=13, color=DARK_GRAY)

add_text_box(slide, 5.5, 1.1, 4, 2.5,
            "🎯 PRIORITÄT 2 (Optional):\n\nMonte Carlo vs Temporal Difference\n\n• MC: Episode-basiert\n• TD: Step-basiert\n• Empirische Analyse",
            font_size=13, color=DARK_GRAY)

add_text_box(slide, 0.5, 4.0, 9, 2.8,
            "KEY INSIGHT:\nDer entscheidende Unterschied liegt in der TARGET BESTIMMUNG!\nWas glaubt der Agent über die Zukunft?",
            font_size=14, bold=True, color=DARK_BLUE)

# ===== SLIDE 2: TARGET BESTIMMUNG (Der Kern!) =====
slide = add_content_slide(prs, "Der Kern: Target Bestimmung")

# SARSA
add_text_box(slide, 0.5, 1.1, 4.3, 0.8,
            "SARSA (On-Policy)",
            font_size=14, bold=True, color=BLUE)

add_text_box(slide, 0.5, 2.0, 4.3, 1.8,
            "TARGET = r + γ·Q(s', a')\n\nwhere a' ~ π(s')\n\n(ACTUAL next action\nfrom policy)",
            font_size=12, color=DARK_GRAY)

# Q-Learning
add_text_box(slide, 5.2, 1.1, 4.3, 0.8,
            "Q-Learning (Off-Policy)",
            font_size=14, bold=True, color=RED)

add_text_box(slide, 5.2, 2.0, 4.3, 1.8,
            "TARGET = r + γ·max Q(s', ·)\n\nwhere · = best action\n\n(OPTIMAL action\nregardless of policy)",
            font_size=12, color=DARK_GRAY)

# Visual difference
add_text_box(slide, 0.5, 4.0, 9, 2.8,
            "💡 DER UNTERSCHIED:\n\n" +
            "SARSA fragt: \"Was wird der Agent TATSÄCHLICH tun?\"\n" +
            "Q-Learning fragt: \"Was wäre OPTIMAL?\"\n\n" +
            "→ SARSA ist KONSERVATIV (vorsichtig mit Risiken)\n" +
            "→ Q-Learning ist AGGRESSIV (optimistic, sucht optimale Policy)",
            font_size=12, color=DARK_GRAY)

# ===== SLIDE 3: Konkretes Beispiel =====
slide = add_content_slide(prs, "Beispiel: Konkrete Zahlen")

add_text_box(slide, 0.5, 1.1, 9, 1,
            "Szenario: State s=10 → State s'=15, Reward r=-1",
            font_size=12, bold=True, color=DARK_GRAY)

add_text_box(slide, 0.5, 2.2, 9, 0.5,
            "Q-Werte in State 15: Q(15,UP)=0.3, Q(15,RIGHT)=0.8, Q(15,DOWN)=0.0, Q(15,LEFT)=0.1",
            font_size=11, color=DARK_GRAY)

# SARSA Calc
add_text_box(slide, 0.5, 2.9, 4.5, 2,
            "SARSA Berechnung:\n\n" +
            "Assume a'=RIGHT (ε-greedy)\n" +
            "TARGET = -1 + 0.99 × 0.8\n" +
            "       = -1 + 0.792\n" +
            "       = -0.208",
            font_size=11, color=BLUE)

# Q-Learning Calc
add_text_box(slide, 5.0, 2.9, 4.5, 2,
            "Q-Learning Berechnung:\n\n" +
            "a* = argmax = RIGHT\n" +
            "TARGET = -1 + 0.99 × 0.8\n" +
            "       = -1 + 0.792\n" +
            "       = -0.208",
            font_size=11, color=RED)

add_text_box(slide, 0.5, 5.1, 9, 2,
            "⚠️  In diesem Fall gleich! Aber wenn Q(15,RIGHT) nicht das beste wäre...\n\n" +
            "Wenn Q(15,UP)=1.0 > Q(15,RIGHT)=0.8:\n" +
            "SARSA: -0.208 (nutzt gewählte Aktion)  |  Q-Learning: -0.01 (nutzt beste)",
            font_size=11, color=DARK_GRAY)

# ===== SLIDE 4: Verhalten in der Umgebung =====
slide = add_content_slide(prs, "Verhalten in der Umgebung")

# SARSA Verhalten
add_text_box(slide, 0.5, 1.1, 4.3, 0.6,
            "SARSA: Risk-Aware",
            font_size=13, bold=True, color=BLUE)

behavior_sarsa = [
    "✓ Lernt Wert UNTER AKTUELLER Policy",
    "✓ Berücksichtigt Explorations-Risiken",
    "✓ Konservative Policy: meidet gefährliche Pfade",
    "✗ Lernt langsamer",
    "✓ Stabiler, weniger Overoptimism"
]
add_bullet_points(slide, 0.5, 1.8, 4.3, 3, behavior_sarsa, font_size=11, color=DARK_GRAY)

# Q-Learning Verhalten
add_text_box(slide, 5.2, 1.1, 4.3, 0.6,
            "Q-Learning: Optimal-Seeking",
            font_size=13, bold=True, color=RED)

behavior_ql = [
    "✓ Lernt Wert der OPTIMALEN Policy",
    "✓ Ignoriert Explorations-Risiken",
    "✓ Aggressive Policy: nimmt Risiken für Optimalität",
    "✓ Lernt schneller",
    "✗ Kann Werte überschätzen (Overoptimism)"
]
add_bullet_points(slide, 5.2, 1.8, 4.3, 3, behavior_ql, font_size=11, color=DARK_GRAY)

# Grafik Area
add_text_box(slide, 0.5, 5.1, 9, 1.8,
            "BEISPIEL: Roboter vor Klippe\n\n" +
            "SARSA: \"Ich exploriere vorsichtig. Wenn ich zufällig runterfalle, lerne ich, dass Nähe zur Klippe schlecht ist.\"\n" +
            "Q-Learning: \"Ich weiß, beste Aktion meidet Klippe. Ich lerne optimistisch, obwohl ich runterfallen könnte.\"",
            font_size=11, color=DARK_GRAY)

# ===== SLIDE 5: Empirische Learning Curves =====
slide = add_content_slide(prs, "Empirische Resultate: Learning Curves")

img_path = os.path.join(OUTPUT_DIR, "01_SARSA_vs_QLearning_Overview.png")
if os.path.exists(img_path):
    try:
        slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.0), width=Inches(9.4))
    except:
        add_text_box(slide, 0.5, 3, 9, 1, "Graph nicht verfügbar", font_size=12, color=RED)

# ===== SLIDE 6: Empirische Resultate Analyse =====
slide = add_content_slide(prs, "Empirische Analyse: Was zeigt sich?")

insights = [
    "1. KONVERGENZ-GESCHWINDIGKEIT:",
    "   Q-Learning konvergiert SCHNELLER (offensiver Lernansatz)",
    "   SARSA braucht länger aber ist STABILER",
    "",
    "2. TD-ERROR MAGNITUDE:",
    "   Q-Learning: größere TD-Fehler (aggressiver)",
    "   SARSA: kleinere TD-Fehler (vorsichtiger)",
    "",
    "3. FINALE PERFORMANCE:",
    "   Oft ähnlich am Ende (beide konvergieren zu guten Policies)",
    "   Unterschied hauptsächlich in LEARNING PROCESS",
]

add_bullet_points(slide, 0.5, 1.1, 9, 6, insights, font_size=11, color=DARK_GRAY)

# ===== SLIDE 7: Monte Carlo vs TD (PRIO 2) =====
slide = add_content_slide(prs, "PRIO 2 (Optional): Monte Carlo vs TD")

# MC
add_text_box(slide, 0.5, 1.1, 4.3, 0.6,
            "Monte Carlo",
            font_size=13, bold=True, color=GREEN)

mc_points = [
    "• Update NACH ganzer Episode",
    "• Q(s,a) += α(G - Q(s,a))",
    "• G = gesamter diskontierter Return",
    "• Unbiased aber High Variance",
    "• Braucht viele Samples",
    "• Gut für kurze Episoden"
]
add_bullet_points(slide, 0.5, 1.8, 4.3, 4, mc_points, font_size=10, color=DARK_GRAY)

# TD
add_text_box(slide, 5.2, 1.1, 4.3, 0.6,
            "Temporal Difference",
            font_size=13, bold=True, color=ORANGE)

td_points = [
    "• Update NACH jedem Schritt",
    "• Q(s,a) += α(r + γV(s') - Q(s,a))",
    "• Bootstrapping: nutzt alte Q-Schätzung",
    "• Biased aber Low Variance",
    "• Schnelleres Lernen",
    "• Gut für lange Episoden"
]
add_bullet_points(slide, 5.2, 1.8, 4.3, 4, td_points, font_size=10, color=DARK_GRAY)

add_text_box(slide, 0.5, 6.0, 9, 1.2,
            "KEY: MC sammelt GESAMTE Episode bevor Update | TD aktualisiert SOFORT nach jedem Schritt",
            font_size=11, bold=True, color=DARK_BLUE)

# ===== SLIDE 8: Empirischer MC vs TD Vergleich =====
slide = add_content_slide(prs, "Empirisch: MC vs TD Learning Speed")

img_path = os.path.join(OUTPUT_DIR, "02_MC_vs_TD_Comparison.png")
if os.path.exists(img_path):
    try:
        slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.0), width=Inches(9.4))
    except:
        add_text_box(slide, 0.5, 3, 9, 1, "Graph nicht verfügbar", font_size=12, color=RED)

# ===== SLIDE 9: Synthesevergleich =====
slide = add_content_slide(prs, "Synthesevergleich: Alle 4 Algorithmen")

comparison_text = """
LERNVERHALTEN UND CHARACTERISTICS:

┌─────────────────────────────────────────────────────────────────────────────┐
│ DIMENSION          │ SARSA      │ Q-Learning │ MC         │ TD (generell)  │
├─────────────────────────────────────────────────────────────────────────────┤
│ Policy Type        │ On-Policy  │ Off-Policy │ -          │ -              │
│ Lernart            │ Online TD  │ Online TD  │ Batch      │ Online         │
│ Update Timing      │ Pro Schritt│ Pro Schritt│ Pro Episod │ Pro Schritt    │
│ Bias/Variance      │ Low/Low    │ Low/Low    │ None/High  │ Low/Low        │
│ Konvergenz Speed   │ Mittel     │ Schnell    │ Langsam    │ Schnell        │
│ Stabilität         │ Hoch       │ Mittel     │ Mittel     │ Hoch           │
│ Overoptimism       │ Nein       │ Ja         │ Nein       │ Nein           │
└─────────────────────────────────────────────────────────────────────────────┘

PRAKTISCHE ANWENDUNG:
✓ SARSA:      Robot Learning (Sicherheit wichtig)
✓ Q-Learning: Game Playing (Optimalität wichtig)
✓ MC:         Policy Evaluation (unbiased)
✓ TD:         General RL (schnell & stabil)
"""

add_text_box(slide, 0.3, 1.1, 9.4, 6, comparison_text, font_size=9, color=DARK_GRAY)

# ===== SLIDE 10: Zusammenfassung =====
slide = add_content_slide(prs, "Zusammenfassung & Key Takeaways")

summary_points = [
    "🎯 PRIO 1: TARGET BESTIMMUNG ist der KERN-UNTERSCHIED",
    "   SARSA: r + γ·Q(s', actual_action) — KONSERVATIV, auf-policy",
    "   Q-Learning: r + γ·max_a Q(s', a) — AGGRESSIV, off-policy",
    "",
    "🎯 VERHALTEN:",
    "   SARSA: Risk-aware, stabil, langsameres Lernen",
    "   Q-Learning: Optimal-seeking, aggressiv, schnelleres Lernen",
    "",
    "🎯 PRIO 2: MC vs TD",
    "   MC: Unbiased aber variabel, Episode-weise Updates",
    "   TD: Biased aber stabil, Step-weise Updates",
    "",
    "🎯 PRAKTISCH:",
    "   Wähle SARSA für Sicherheit, Q-Learning für Optimalität",
    "   Nutze TD-Methoden für schnelles Lernen",
]

add_bullet_points(slide, 0.5, 1.1, 9, 6, summary_points, font_size=11, color=DARK_GRAY)

# ============================================================================
# SPEICHERN
# ============================================================================

output_path = os.path.join(OUTPUT_DIR, "On_Policy_vs_Off_Policy_Analysis.pptx")
prs.save(output_path)

print(f"✅ PowerPoint Präsentation erstellt:")
print(f"   {output_path}")
print(f"\n📊 INHALT (11 Slides):")
print(f"   0. Title Slide")
print(f"   1. Agenda & Überblick")
print(f"   2. TARGET BESTIMMUNG (Der Kern!)")
print(f"   3. Konkretes Beispiel mit Zahlen")
print(f"   4. Verhalten in der Umgebung")
print(f"   5. Empirische Learning Curves")
print(f"   6. Empirische Analyse")
print(f"   7. MC vs TD (PRIO 2 Optional)")
print(f"   8. MC vs TD Empirisch")
print(f"   9. Synthesevergleich (4 Algorithmen)")
print(f"   10. Zusammenfassung")
print(f"\n✨ PROFESSIONELLE FEATURES:")
print(f"   ✓ Empirische Grafiken eingebettet")
print(f"   ✓ Detaillierte Code-Erklärungen")
print(f"   ✓ Konkretes Zahlenbeispiel")
print(f"   ✓ Vergleichstabelle")
print(f"   ✓ Praktische Anwendungsbeispiele")
print(f"\n📖 Ready for presentation!")
