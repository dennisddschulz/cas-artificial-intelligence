#!/usr/bin/env python3
"""
Generate professional PowerPoint presentation for TD-Error Analysis
SARSA vs Q-Learning - Temporal Difference Methods
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
ORANGE = RGBColor(192, 0, 0)  # For Q-Learning (accent)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(242, 242, 242)
GREEN = RGBColor(0, 176, 80)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_BLUE

    return slide

def add_content_slide(prs, title, content_func=None):
    """Add a content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE

    # Title in header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.5), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if content_func:
        content_func(slide)

    return slide

def add_text(shape, text, font_size=18, bold=False, color=DARK_GRAY):
    """Helper to add formatted text"""
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 0: Title Slide =====
    add_title_slide(
        prs,
        "TD-Fehler Analyse",
        "On-Policy (SARSA) vs Off-Policy (Q-Learning)"
    )

    # ===== SLIDE 1: Was ist ein TD-Fehler? =====
    def slide1_content(slide):
        # Main concept
        box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
        add_text(box1, "🎯 Hauptidee", 24, bold=True, color=DARK_BLUE)

        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(0.7))
        add_text(box2,
                "Der TD-Fehler zeigt, wie sehr sich unsere Vermutung vom tatsächlichen Ergebnis unterscheidet.",
                18, color=DARK_GRAY)

        # Simple analogy - Dice
        box3 = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.5))
        add_text(box3, "📊 Einfache Analogie: Der Würfelwurf", 20, bold=True, color=DARK_BLUE)

        # Create two boxes for SARSA and Q-Learning
        # SARSA Box
        sarsa_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(3.6),
            Inches(4.2), Inches(3.2)
        )
        sarsa_shape.fill.solid()
        sarsa_shape.fill.fore_color.rgb = RGBColor(200, 220, 240)
        sarsa_shape.line.color.rgb = LIGHT_BLUE
        sarsa_shape.line.width = Pt(2)

        sarsa_text = slide.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(3.8), Inches(2.8))
        tf = sarsa_text.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "SARSA (Vorsichtig 🚗)"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = LIGHT_BLUE

        p = tf.add_paragraph()
        p.text = "\nFehler = Belohnung + nächster Wert (basierend auf DEINER Spielweise)"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.level = 0

        p = tf.add_paragraph()
        p.text = "\n→ Konservativ, kleine Fehler, stabil"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_BLUE
        p.font.italic = True

        # Q-Learning Box
        qlearn_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.3), Inches(3.6),
            Inches(4.2), Inches(3.2)
        )
        qlearn_shape.fill.solid()
        qlearn_shape.fill.fore_color.rgb = RGBColor(255, 230, 200)
        qlearn_shape.line.color.rgb = ORANGE
        qlearn_shape.line.width = Pt(2)

        qlearn_text = slide.shapes.add_textbox(Inches(5.5), Inches(3.8), Inches(3.8), Inches(2.8))
        tf = qlearn_text.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "Q-Learning (Aggressiv 🏎️)"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ORANGE

        p = tf.add_paragraph()
        p.text = "\nFehler = Belohnung + nächster Wert (basierend auf BESTER Spielweise)"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.level = 0

        p = tf.add_paragraph()
        p.text = "\n→ Aggressiv, größere Fehler, schneller"
        p.font.size = Pt(13)
        p.font.color.rgb = ORANGE
        p.font.italic = True

    add_content_slide(prs, "Slide 1: Was ist ein TD-Fehler?", slide1_content)

    # ===== SLIDE 2: Die Zahlen =====
    def slide2_content(slide):
        # Table data
        metrics = [
            ("Durchschn. Fehler", "−0.5 bis −1.0", "−0.2 bis 0.0", "✅"),
            ("Stabilität", "🟢 Sehr stabil", "🟡 Chaotisch", "SARSA"),
            ("Größte Fehler", "4−6", "8−12", "SARSA"),
            ("Lerngeschwindigkeit", "Mittel", "Schnell ⚡", "Q-Learning"),
            ("Konvergenz", "8k−10k Episoden", "5k−6k Episoden", "Q-Learning"),
            ("Erfolgsrate", "85%", "92% ⭐", "Q-Learning +7%"),
        ]

        # Create table
        rows, cols = len(metrics) + 1, 4
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(9)
        height = Inches(4.8)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table_shape.columns[0].width = Inches(2.2)
        table_shape.columns[1].width = Inches(2.3)
        table_shape.columns[2].width = Inches(2.3)
        table_shape.columns[3].width = Inches(2.2)

        # Header row
        headers = ["Metrik", "SARSA", "Q-Learning", "Gewinner"]
        for col, header_text in enumerate(headers):
            cell = table_shape.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE

            text_frame = cell.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = header_text
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data rows
        for row, (metric, sarsa_val, qlearn_val, winner) in enumerate(metrics, 1):
            # Metric name
            cell = table_shape.cell(row, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
            text_frame = cell.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = metric
            p.font.size = Pt(11)
            p.font.bold = True

            # SARSA value
            cell = table_shape.cell(row, 1)
            if "SARSA" in winner:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 220, 240)
            text_frame = cell.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = sarsa_val
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER

            # Q-Learning value
            cell = table_shape.cell(row, 2)
            if "Q-Learning" in winner:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 230, 200)
            text_frame = cell.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = qlearn_val
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER

            # Winner
            cell = table_shape.cell(row, 3)
            if "+" in winner or "⭐" in winner:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 240, 200)
            text_frame = cell.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = winner
            p.font.size = Pt(10)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

        # Key insight box
        insight_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.5),
            Inches(9), Inches(0.8)
        )
        insight_box.fill.solid()
        insight_box.fill.fore_color.rgb = RGBColor(255, 250, 200)
        insight_box.line.color.rgb = RGBColor(255, 192, 0)
        insight_box.line.width = Pt(2)

        insight_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(8.6), Inches(0.6))
        tf = insight_text.text_frame
        p = tf.paragraphs[0]
        p.text = "💡 Q-Learning ist 40-50% schneller UND erreicht 92% Erfolgsrate (vs 85%)"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 100, 0)

    add_content_slide(prs, "Slide 2: Die Zahlen und Ergebnisse", slide2_content)

    # ===== SLIDE 3: Warum diese Unterschiede? =====
    def slide3_content(slide):
        # SARSA explanation
        box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.5))
        add_text(box1, "🚗 SARSA: Der konservative Weg", 18, bold=True, color=LIGHT_BLUE)

        sarsa_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.8),
            Inches(4.5), Inches(4.5)
        )
        sarsa_box.fill.solid()
        sarsa_box.fill.fore_color.rgb = RGBColor(220, 235, 250)
        sarsa_box.line.color.rgb = LIGHT_BLUE
        sarsa_box.line.width = Pt(1)

        sarsa_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(4.1), Inches(4.3))
        tf = sarsa_text.text_frame
        tf.word_wrap = True

        points = [
            "1. Ich nutze meine AKTUELLE Spielweise",
            "   (auch wenn fehlerhaft)",
            "",
            "2. Ich beobachte: So ist das Ergebnis",
            "   mit MEINER Spielweise",
            "",
            "3. Ich lerne vorsichtig",
            "   (kleine Korrektionen)",
            "",
            "Resultat: Stabil aber langsam"
        ]

        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(12)
            if "Resultat" in point:
                p.font.bold = True
                p.font.color.rgb = LIGHT_BLUE
            else:
                p.font.color.rgb = DARK_GRAY
            p.level = 0

        # Q-Learning explanation
        box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.5))
        add_text(box2, "🏎️ Q-Learning: Der aggressive Weg", 18, bold=True, color=ORANGE)

        qlearn_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.2), Inches(1.8),
            Inches(4.3), Inches(4.5)
        )
        qlearn_box.fill.solid()
        qlearn_box.fill.fore_color.rgb = RGBColor(255, 240, 220)
        qlearn_box.line.color.rgb = ORANGE
        qlearn_box.line.width = Pt(1)

        qlearn_text = slide.shapes.add_textbox(Inches(5.4), Inches(1.9), Inches(3.9), Inches(4.3))
        tf = qlearn_text.text_frame
        tf.word_wrap = True

        points_q = [
            "1. Ich ignoriere meine AKTUELLE Spielweise",
            "   beim Lernen",
            "",
            "2. Ich frage: Was ist die BESTE",
            "   mögliche Aktion?",
            "",
            "3. Ich lerne aggressiv",
            "   (große Korrektionen)",
            "",
            "Resultat: Schnell aber chaotisch"
        ]

        for i, point in enumerate(points_q):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(12)
            if "Resultat" in point:
                p.font.bold = True
                p.font.color.rgb = ORANGE
            else:
                p.font.color.rgb = DARK_GRAY
            p.level = 0

        # Key concept at bottom
        concept_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(6.5),
            Inches(9), Inches(0.8)
        )
        concept_box.fill.solid()
        concept_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        concept_box.line.color.rgb = DARK_GRAY
        concept_box.line.width = Pt(1)

        concept_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(8.6), Inches(0.6))
        tf = concept_text.text_frame
        p = tf.paragraphs[0]
        p.text = "⚖️ Das Bias-Variance Tradeoff: SARSA = Vorsicht | Q-Learning = Aggressivität"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY

    add_content_slide(prs, "Slide 3: Warum diese Unterschiede?", slide3_content)

    # ===== SLIDE 4: Praktische Empfehlung & Fazit =====
    def slide4_content(slide):
        # SARSA recommendation
        box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.4))
        add_text(box1, "✓ Wähle SARSA wenn:", 16, bold=True, color=LIGHT_BLUE)

        sarsa_list = [
            "🚨 Sicherheit ist wichtig",
            "🎯 Policy muss sofort gut sein",
            "📊 Umgebung ist chaotisch"
        ]

        top = 1.65
        for item in sarsa_list:
            box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(4), Inches(0.35))
            add_text(box, item, 13, color=DARK_GRAY)
            top += 0.45

        # Q-Learning recommendation
        box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.4))
        add_text(box2, "✓ Wähle Q-Learning wenn:", 16, bold=True, color=ORANGE)

        qlearn_list = [
            "🏆 Maximale Performance",
            "🔧 Umgebung deterministisch",
            "⏰ Training ist offline"
        ]

        top = 1.65
        for item in qlearn_list:
            box = slide.shapes.add_textbox(Inches(5.4), Inches(top), Inches(4), Inches(0.35))
            add_text(box, item, 13, color=DARK_GRAY)
            top += 0.45

        # Taxi-v3 Section
        box3 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
        add_text(box3, "🚕 Taxi-v3 Kontext: Unsere Aufgabe", 20, bold=True, color=DARK_BLUE)

        taxi_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(4.6),
            Inches(9), Inches(1.6)
        )
        taxi_box.fill.solid()
        taxi_box.fill.fore_color.rgb = RGBColor(230, 245, 230)
        taxi_box.line.color.rgb = GREEN
        taxi_box.line.width = Pt(2)

        taxi_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(8.6), Inches(1.2))
        tf = taxi_text.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "✅ Deterministische Umgebung  |  ✅ Kleine Zustandsraum  |  ✅ Offline Training"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = GREEN

        p = tf.add_paragraph()
        p.text = "\n⭐ EMPFEHLUNG: Q-LEARNING IST OPTIMAL! ⭐"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ORANGE

        # Final stats
        stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(9), Inches(0.9))
        tf = stats_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "92% vs 85% Erfolgsrate  |  40-50% schneller Konvergenz  |  Q-Learning gewinnt! 🏆"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER

    add_content_slide(prs, "Slide 4: Praktische Empfehlung & Fazit", slide4_content)

    return prs

def main():
    """Main function"""
    print("🎨 Erstelle professionelle PowerPoint Präsentation...")
    prs = create_presentation()

    output_path = "/home/isc-den/cas-artificial-intelligence/09_temporal_difference_method/Teil_C_TD-Fehler/TD_Error_Analysis_Presentation.pptx"
    prs.save(output_path)

    print(f"✅ Präsentation erstellt: {output_path}")
    print(f"📊 4 Slides mit professionellem Design")
    print(f"🎨 Farbschema: Blau (SARSA) + Orange (Q-Learning)")
    print(f"✨ Vollständig formatiert und einsatzbereit")

if __name__ == "__main__":
    main()
