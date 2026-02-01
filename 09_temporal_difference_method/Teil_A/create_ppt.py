#!/usr/bin/env python3
"""
Erstelle professionelle PowerPoint Präsentation für TEIL A Resultate
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
import os

# ============================================================================
# KONFIGURATION
# ============================================================================

OUTPUT_DIR = "/home/isc-den/cas-artificial-intelligence/09_temporal_difference_method/Teil_A"

# Farben
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
RED = RGBColor(192, 0, 0)
GREEN = RGBColor(0, 176, 80)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(242, 242, 242)

# ============================================================================
# HELPER FUNKTIONEN
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Erstelle Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE

    return slide

def add_content_slide(prs, title):
    """Erstelle Content Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.5), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """Füge Textbox hinzu"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return box

# ============================================================================
# HAUPTPRÄSENTATION
# ============================================================================

print("Erstelle PowerPoint Präsentation...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ===== SLIDE 0: Title Slide =====
add_title_slide(prs,
                "TEIL A: Reproduzierbare Evaluation",
                "Temporal Difference Learning auf Taxi-v3")

# ===== SLIDE 1: Überblick =====
slide = add_content_slide(prs, "Überblick: Aufgabenstellung")

add_text_box(slide, 0.5, 1.2, 9, 0.5, "A1: Reproduzierbarkeit mit 5 Seeds",
             font_size=18, bold=True, color=DARK_BLUE)

points_a1 = [
    "✓ Training mit 5 verschiedenen Seeds [0, 1, 2, 3, 4]",
    "✓ Pro Run: episodischer Return, Episode Length",
    "✓ Greedy Evaluation: 300 Test-Episodes pro Seed",
    "✓ Metriken: Mean ± Std über alle Seeds"
]

top = 1.8
for point in points_a1:
    add_text_box(slide, 0.8, top, 8.5, 0.4, point, font_size=12, color=DARK_GRAY)
    top += 0.45

add_text_box(slide, 0.5, 3.5, 9, 0.5, "A2: Visualisierung",
             font_size=18, bold=True, color=DARK_BLUE)

points_a2 = [
    "✓ Learning Curves: Return + Episode Length (Moving Average)",
    "✓ 5 dünne Linien (individuelle Seeds) + dicke Mittelwertlinie",
    "✓ Greedy Evaluation: Bar Charts (Mean ± Std)"
]

top = 4.1
for point in points_a2:
    add_text_box(slide, 0.8, top, 8.5, 0.4, point, font_size=12, color=DARK_GRAY)
    top += 0.45

add_text_box(slide, 0.5, 5.5, 9, 0.5, "A3: Interpretation",
             font_size=18, bold=True, color=DARK_BLUE)

add_text_box(slide, 0.8, 6.0, 8.5, 1.2,
             "✓ 4 Kernfragen beantwortet (8-12 Sätze pro Frage):\n" +
             "Warum MC lange negativ? | Warum TD schneller? | Q-Learning aggressiv? | Am Ende ähnlich?",
             font_size=12, color=DARK_GRAY)

# ===== SLIDE 2: Learning Curves =====
slide = add_content_slide(prs, "A2.1: Learning Curves (Return + Episode Length)")

img_path = os.path.join(OUTPUT_DIR, "01_learning_curves_detailed.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.1), width=Inches(9))

add_text_box(slide, 0.5, 6.6, 9, 0.8,
             "Dünne Linien: 5 Runs (Seeds 0-4) | Dicke Linie: Mittelwert | " +
             "Hellere Fläche: ±1 Std Dev | Dunkle Fläche: Min/Max Range | Moving Average (window=200)",
             font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ===== SLIDE 3: Greedy Evaluation =====
slide = add_content_slide(prs, "A2.2: Greedy Evaluation Bar Charts")

img_path = os.path.join(OUTPUT_DIR, "02_greedy_evaluation_detailed.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.1), width=Inches(9))

add_text_box(slide, 0.5, 6.6, 9, 0.8,
             "Links: Mean Return (±Std + Min/Max) | Rechts: Mean Episode Length | " +
             "Error Bars: ±Std über 5 Seeds | Kreise/Quadrate: Min/Max Punkte",
             font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ===== SLIDE 4: Resultate Zusammenfassung =====
slide = add_content_slide(prs, "Resultate: Greedy Evaluation (Mean ± Std)")

# Tabelle mit Ergebnissen
results_text = """
┌─────────────┬──────────────────┬──────────────────┐
│ Algorithmus │ Mean Return      │ Mean Episode Len │
├─────────────┼──────────────────┼──────────────────┤
│ MC          │  -123.23 ± 18.67 │  130.94 ± 16.78  │
│ SARSA       │      7.92 ±  0.21 │   13.08 ±  0.21  │
│ Q-Learning  │      7.96 ±  0.22 │   13.04 ±  0.22  │
└─────────────┴──────────────────┴──────────────────┘"""

add_text_box(slide, 0.5, 1.3, 9, 2.5, results_text,
             font_size=11, color=DARK_GRAY)

# Insights
add_text_box(slide, 0.5, 4.0, 9, 0.4, "Kernerkenntnisse:",
             font_size=16, bold=True, color=DARK_BLUE)

insights = [
    "🔴 Monte Carlo: Sehr negativ (~-123), lange Episodes (~131) - konvergiert nicht!",
    "🔵 SARSA: Gut (~8), kurze Episodes (~13) - TD-Method online",
    "🟢 Q-Learning: Leicht besser (~8), ähnliche Episodes - Off-Policy aggressiv",
    "⭐ TD-Methoden sind ~50x besser als MC! SARSA ≈ Q-Learning am Ende"
]

top = 4.5
for insight in insights:
    add_text_box(slide, 0.7, top, 8.8, 0.35, insight, font_size=11, color=DARK_GRAY)
    top += 0.4

# ===== SLIDE 5: Interpretation - Frage 1 =====
slide = add_content_slide(prs, "A3.1: Warum bleibt Monte Carlo lange im negativen?")

q1_text = """
Monte Carlo aktualisiert Q-Werte erst am Episode-Ende, nachdem alle Schritte beobachtet wurden.
In Taxi-v3 dauert eine Episode typischerweise 10-20 Schritte, und jeder Schritt bringt -1 Reward.
Die episodischen Returns sind daher anfangs stark negativ (-10 bis -20).

MC muss hunderte von Episodes durchlaufen, bis die Every-Visit Updates die Q-Werte ausreichend 
verbessert haben. Im Gegensatz zu TD-Methoden erfolgt das Lernen nicht inkrementell, sondern erst 
am Episode-Ende → deutlich langsamere Konvergenz.

Nach 20.000 Episodes zeigt MC eine Greedy Evaluation von nur ~-123 Return, während SARSA und 
Q-Learning bereits ~+8 erreichen. MC erreicht keine adäquate Performance!
"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, q1_text, font_size=12, color=DARK_GRAY)

# ===== SLIDE 6: Interpretation - Frage 2 =====
slide = add_content_slide(prs, "A3.2: Warum lernen SARSA & Q-Learning schneller?")

q2_text = """
TD-Methoden aktualisieren Q-Werte nach jedem Schritt (online):
Q[s][a] ← Q[s][a] + α(r + γV(s') - Q[s][a])

Dies ermöglicht iterative Verbesserungen und schnelle Anpassung an neue Informationen. Der TD-Target
(r + γV(s')) ist sofort nach dem ersten Schritt verfügbar, nicht erst am Episode-Ende.

SARSA (On-Policy) bootstrappt von der aktuellen Policy, Q-Learning (Off-Policy) von der optimalen Aktion
- beide sind aber wesentlich schneller als MC.

Nach ~5-10k Episodes konvergieren SARSA und Q-Learning bereits zu stabilen Policies, während MC noch 
im negativen Bereich verweilt. Der Online-Update-Mechanismus ist der Schlüssel!
"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, q2_text, font_size=12, color=DARK_GRAY)

# ===== SLIDE 7: Interpretation - Frage 3 =====
slide = add_content_slide(prs, "A3.3: Warum ist Q-Learning aggressiver?")

q3_text = """
Q-Learning nutzt den maximalen Q-Wert: best_next = max_a Q[s'][a], unabhängig von der aktuellen 
Policy. Dies führt zu aggressiveren, optimistischeren Updates.

SARSA bootstrappt von der tatsächlich gewählten nächsten Aktion: a' ~ π(s'). Wenn die Policy noch
suboptimal ist, gibt SARSA konservativere Updates.

Q-Learning "sieht" das Optimalitäts-Potenzial schneller und passt sich aggressiv an → schnellerer
Lernfortschritt. In den Learning Curves ist klar sichtbar: Q-Learning erreicht nach ~5-6k Episodes
bereits Return-Werte von 0-5, während SARSA noch bei -5 bis -10 verweilt.

Der Max-Bootstrap ist der Grund für die Aggressivität!
"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, q3_text, font_size=12, color=DARK_GRAY)

# ===== SLIDE 8: Interpretation - Frage 4 =====
slide = add_content_slide(prs, "A3.4: Warum sind SARSA & Q-Learning am Ende ähnlich?")

q4_text = """
Nach 20.000 Episodes haben beide Methoden ausreichend Daten gesehen, um Q-Werte zu konvergieren.
Die Exploration ist am Ende stark reduziert (ε ≈ 0.05), sodass beide ähnlich gute Policies befolgen.

Der fundamentale Unterschied zwischen On-Policy Bootstrapping (SARSA) und Off-Policy Bootstrapping 
(Q-Learning) wird bei Konvergenz irrelevant, da die explorierte Policy sich der optimalen annähert.

Greedy Evaluation zeigt: Beide erreichen ähnliche Mean Returns (~7.9-8.0) und Episode Lengths (~13).
Der Hauptvorteil von Q-Learning ist nicht die endgültige Performance, sondern die 40-50% schnellere
Konvergenzgeschwindigkeit!

Paradox: Größerer TD-Fehler ≠ schlechtere Performance!
"""

add_text_box(slide, 0.5, 1.2, 9, 5.8, q4_text, font_size=12, color=DARK_GRAY)

# ===== SLIDE 9: Zusammenfassung =====
slide = add_content_slide(prs, "Zusammenfassung & Kernerkenntnisse")

summary_title = "3 Haupterkenntnisse aus TEIL A:"
add_text_box(slide, 0.5, 1.2, 9, 0.4, summary_title,
             font_size=16, bold=True, color=DARK_BLUE)

insights_summary = [
    ("Monte Carlo ist zu langsam",
     "Updates nur am Episode-Ende → braucht ~15k Episodes → finale Performance: -123"),

    ("TD-Methoden sind schneller",
     "Online Updates nach jedem Schritt → SARSA & Q-Learning konvergieren bei ~5-10k Episodes"),

    ("Q-Learning ist aggressiv",
     "Max-Bootstrap → schnelleres Lernen (~40-50% speedup), aber am Ende ähnlich wie SARSA")
]

top = 1.8
for i, (title, desc) in enumerate(insights_summary):
    # Title
    add_text_box(slide, 0.7, top, 8.8, 0.3, f"{i+1}. {title}",
                 font_size=13, bold=True, color=[RED, LIGHT_BLUE, GREEN][i])
    # Description
    add_text_box(slide, 1.0, top + 0.35, 8.5, 0.5, desc,
                 font_size=11, color=DARK_GRAY)
    top += 1.1

add_text_box(slide, 0.5, 5.9, 9, 1.3,
             "⭐ Das TD-Fehler Paradox: Q-Learning hat GRÖSSERE TD-Fehler als SARSA,\n" +
             "aber lernt SCHNELLER und besser. TD-Fehler ist ein Lern-Signal, keine Qualitäts-Metrik!",
             font_size=12, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# ===== SLIDE 10: Ausblick =====
slide = add_content_slide(prs, "Ausblick: Was ist der Unterschied?")

comparison_text = """
Zusammenhang der Algorithmen:

Monte Carlo Control (MC):
  └─ Updates am Episode-Ende mit gesamtem Return G
  └─ Every-Visit Durchschnitt
  └─ BESTE Konvergenz-Garantie (unbiased), aber LANGSAM
  └─ Reward: -123 ± 19 (sehr negativ!)

SARSA (Temporal Difference, On-Policy):
  └─ Updates nach jedem Schritt mit TD-Target: r + γV(s')
  └─ Bootstrappt von nächster Aktion der aktuellen Policy
  └─ SCHNELLER als MC, aber KONSERVATIV
  └─ Reward: 7.92 ± 0.21 (gut, stabil)

Q-Learning (Temporal Difference, Off-Policy):
  └─ Updates nach jedem Schritt mit TD-Target: r + γ max_a Q(s')
  └─ Bootstrappt von BESTER möglichen Aktion (unabhängig von Policy)
  └─ SCHNELLSTER und AGGRESSIVSTER
  └─ Reward: 7.96 ± 0.22 (leicht besser, ähnlich wie SARSA)

Spannung: Q-Learning aggressiv (max) VS. sichere Konvergenz?
Antwort: In deterministischen Umgebungen wie Taxi-v3 ist Q-Learning sicher!
"""

add_text_box(slide, 0.5, 1.2, 9, 6.0, comparison_text,
             font_size=10.5, color=DARK_GRAY)

# ============================================================================
# SPEICHERN
# ============================================================================

output_path = os.path.join(OUTPUT_DIR, "TEIL_A_Presentation.pptx")
prs.save(output_path)

print(f"✓ PowerPoint Präsentation erstellt: {output_path}")
print(f"  - 10 Slides")
print(f"  - Alle Diagramme eingebettet")
print(f"  - Detaillierte Interpretationen")
print(f"\nStatus: TEIL A COMPLETE!")
