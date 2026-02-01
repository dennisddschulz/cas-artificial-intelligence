#!/usr/bin/env python3
"""
Erstelle finale, hochprofessionelle PowerPoint mit detaillierten Visualisierungen
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================================
# KONFIGURATION
# ============================================================================

OUTPUT_DIR = "/home/isc-den/cas-artificial-intelligence/09_temporal_difference_method/Teil_A"

# Farben
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
RED = RGBColor(192, 0, 0)
GREEN = RGBColor(0, 176, 80)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(242, 242, 242)

# ============================================================================
# HELPER FUNKTIONEN
# ============================================================================

def add_title_slide(prs, title, subtitle):
    """Erstelle Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE

    return slide

def add_content_slide(prs, title):
    """Erstelle Content Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.5), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """Füge Textbox hinzu"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return box

# ============================================================================
# HAUPTPRÄSENTATION
# ============================================================================

print("Erstelle finale professionelle PowerPoint mit detaillierten Grafiken...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ===== SLIDE 0: Title Slide =====
add_title_slide(prs,
                "TEIL A: Reproduzierbare Evaluation",
                "Temporal Difference Learning auf Taxi-v3")

# ===== SLIDE 1: Überblick =====
slide = add_content_slide(prs, "Überblick: Aufgabenstellung TEIL A")

add_text_box(slide, 0.5, 1.2, 9, 0.5, "A1: Reproduzierbarkeit mit 5 Seeds",
             font_size=18, bold=True, color=DARK_BLUE)

points_a1 = [
    "✓ Training mit 5 verschiedenen Seeds [0, 1, 2, 3, 4]",
    "✓ Pro Run: episodischer Return, Episode Length, Greedy Evaluation (300 Episodes)",
    "✓ Metriken: Mean ± Std ± Min/Max über alle 5 Seeds"
]

top = 1.8
for point in points_a1:
    add_text_box(slide, 0.8, top, 8.5, 0.35, point, font_size=12, color=DARK_GRAY)
    top += 0.4

add_text_box(slide, 0.5, 3.3, 9, 0.5, "A2: Detaillierte Visualisierung",
             font_size=18, bold=True, color=DARK_BLUE)

points_a2 = [
    "✓ Learning Curves: Return + Episode Length (Moving Average window=200)",
    "✓ 5 Runs einzeln als dünne Linien, dicke Mittelwertlinie, ±1σ + Min/Max Bereiche",
    "✓ Greedy Evaluation: Bar Charts mit Mean ± Std und Min/Max Anzeige"
]

top = 3.9
for point in points_a2:
    add_text_box(slide, 0.8, top, 8.5, 0.35, point, font_size=12, color=DARK_GRAY)
    top += 0.4

add_text_box(slide, 0.5, 5.4, 9, 0.5, "A3: Interpretation (8-12 Sätze pro Frage)",
             font_size=18, bold=True, color=DARK_BLUE)

add_text_box(slide, 0.8, 5.95, 8.5, 1.3,
             "✓ Q1: Warum bleibt MC lange im negativen?  |  Q2: Warum TD schneller?\n" +
             "✓ Q3: Warum ist Q-Learning aggressiver?  |  Q4: Warum SARSA & QL ähnlich?",
             font_size=12, color=DARK_GRAY)

# ===== SLIDE 2: Learning Curves (detailliert) =====
slide = add_content_slide(prs, "A2.1: Learning Curves mit detaillierter Streuung")

img_path = os.path.join(OUTPUT_DIR, "01_learning_curves_detailed.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(0.95), width=Inches(9.4))
else:
    add_text_box(slide, 0.5, 3, 9, 1, "Grafik nicht gefunden!", font_size=14, color=RED)

# ===== SLIDE 3: Greedy Evaluation (detailliert) =====
slide = add_content_slide(prs, "A2.2: Greedy Evaluation Bar Charts mit Min/Max")

img_path = os.path.join(OUTPUT_DIR, "02_greedy_evaluation_detailed.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(0.95), width=Inches(9.4))
else:
    add_text_box(slide, 0.5, 3, 9, 1, "Grafik nicht gefunden!", font_size=14, color=RED)

# ===== SLIDE 4: Resultate Tabelle =====
slide = add_content_slide(prs, "Resultate: Greedy Evaluation (Mean ± Std, Min/Max)")

results_text = """
┌────────────┬──────────────────────────┬────────────────────────────┐
│ Algoritm   │ Mean Return (±Std)       │ Mean Episode Len (±Std)    │
├────────────┼──────────────────────────┼────────────────────────────┤
│ MC         │ -123.23 ± 18.67          │ 130.94 ± 16.78             │
│            │ [min: -147.45, max: -97] │ [min: 108, max: 153]       │
├────────────┼──────────────────────────┼────────────────────────────┤
│ SARSA      │    7.92 ± 0.21           │  13.08 ± 0.21              │
│            │ [min: 7.73, max: 8.28]   │ [min: 12.72, max: 13.27]   │
├────────────┼──────────────────────────┼────────────────────────────┤
│ Q-Learning │    7.96 ± 0.22           │  13.04 ± 0.22              │
│            │ [min: 7.64, max: 8.28]   │ [min: 12.72, max: 13.36]   │
└────────────┴──────────────────────────┴────────────────────────────┘"""

add_text_box(slide, 0.3, 1.2, 9.4, 3.2, results_text,
             font_size=9.5, color=DARK_GRAY)

add_text_box(slide, 0.5, 4.6, 9, 0.4, "⭐ Kernerkenntnisse aus den Daten:",
             font_size=13, bold=True, color=DARK_BLUE)

insights = [
    "🔴 Monte Carlo: Sehr negativ (-123), lange Episodes (131) - konvergiert nicht!",
    "🔵 SARSA: Gut (7.92), kurze Episodes (13.08) - TD-Methode online",
    "🟢 Q-Learning: Leicht besser (7.96), ähnlich wie SARSA - aggressiv aber konvergent",
    "⭐ TD-Methoden sind ~50x besser als MC! Streuung sehr gering bei SARSA & Q-Learning."
]

top = 5.1
for insight in insights:
    add_text_box(slide, 0.7, top, 8.8, 0.3, insight, font_size=11, color=DARK_GRAY)
    top += 0.35

# ===== SLIDE 5: Interpretation Q1 =====
slide = add_content_slide(prs, "A3.1: Warum bleibt Monte Carlo lange im negativen?")

q1_text = """
Monte Carlo aktualisiert Q-Werte erst am Ende einer Episode, nachdem alle Schritte beobachtet wurden.
In Taxi-v3 dauert eine Episode typischerweise 10-20 Schritte, und jeder Schritt bringt -1 Reward.
Daher sind die episodischen Returns anfangs stark negativ (z.B. -10 bis -20).

MC muss hunderte von Episodes durchlaufen, bis die Every-Visit Updates die Q-Werte ausreichend 
verbessert haben. Im Gegensatz zu TD-Methoden erfolgt das Lernen nicht inkrementell, sondern erst 
am Episode-Ende → deutlich langsamere Konvergenz.

Nach 20.000 Episodes zeigt MC noch immer eine Greedy Evaluation von -123 Return (±19 Streuung),
während SARSA und Q-Learning bereits +8 erreichen. MC konvergiert für diese Aufgabe faktisch nicht!
"""

add_text_box(slide, 0.5, 1.2, 9, 6.0, q1_text, font_size=12, color=DARK_GRAY)

# ===== SLIDE 6: Interpretation Q2 =====
slide = add_content_slide(prs, "A3.2: Warum lernen SARSA & Q-Learning schneller?")

q2_text = """
TD-Methoden aktualisieren Q-Werte nach JEDEM Schritt (online):
  Q[s][a] ← Q[s][a] + α(r + γV(s') - Q[s][a])

Dies ermöglicht iterative Verbesserungen und schnelle Anpassung an neue Informationen. Der TD-Target
(r + γV(s')) ist sofort nach dem ersten Schritt verfügbar, nicht erst am Episode-Ende.

SARSA (On-Policy) bootstrappt von der aktuellen Policy, Q-Learning (Off-Policy) von der optimalen 
Aktion - beide sind aber wesentlich schneller als MC.

Die Ergebnisse zeigen: Nach ~5-10k Episodes konvergieren SARSA und Q-Learning bereits zu stabilen 
Policies (Return ~8±0.2), während MC noch weit entfernt ist. Online-Updates sind der Schlüssel!
"""

add_text_box(slide, 0.5, 1.2, 9, 6.0, q2_text, font_size=12, color=DARK_GRAY)

# ===== SLIDE 7: Interpretation Q3 =====
slide = add_content_slide(prs, "A3.3: Warum ist Q-Learning aggressiver?")

q3_text = """
Q-Learning nutzt den maximalen Q-Wert: best_next = max_a Q[s'][a], unabhängig von der aktuellen 
Policy. Dies führt zu aggressiveren, optimistischeren Updates.

SARSA bootstrappt von der tatsächlich gewählten nächsten Aktion: a' ~ π(s'). Wenn die Policy noch
suboptimal ist, gibt SARSA konservativere Updates.

Q-Learning "sieht" das Optimalitäts-Potenzial schneller und passt sich aggressiv an. In den Learning 
Curves ist klar sichtbar: Q-Learning erreicht nach ~5-6k Episodes bereits Return-Werte von 0-5, 
während SARSA noch bei -5 bis -10 verweilt.

Der Max-Bootstrap ist der Grund für die Aggressivität - und damit auch für die 40-50% schnellere Konvergenz!
"""

add_text_box(slide, 0.5, 1.2, 9, 6.0, q3_text, font_size=12, color=DARK_GRAY)

# ===== SLIDE 8: Interpretation Q4 =====
slide = add_content_slide(prs, "A3.4: Warum sind SARSA & Q-Learning am Ende ähnlich?")

q4_text = """
Nach 20.000 Episodes haben beide Methoden ausreichend Daten gesehen, um Q-Werte zu konvergieren.
Die Exploration ist am Ende stark reduziert (ε ≈ 0.05), sodass beide Algorithmen ähnlich gute 
Policies befolgen.

Der fundamentale Unterschied zwischen On-Policy Bootstrapping (SARSA) und Off-Policy Bootstrapping 
(Q-Learning) wird bei Konvergenz irrelevant, da die explorierte Policy sich der optimalen nähert.

Die Greedy Evaluation zeigt: Beide erreichen ähnliche Mean Returns (~7.9-8.0 ± 0.2) und Episode 
Lengths (~13.0 ± 0.2). Der Hauptvorteil von Q-Learning ist nicht die endgültige Performance, 
sondern die 40-50% schnellere Konvergenzgeschwindigkeit!

Paradox gelöst: Größerer TD-Fehler ≠ schlechtere Performance!
"""

add_text_box(slide, 0.5, 1.2, 9, 6.0, q4_text, font_size=12, color=DARK_GRAY)

# ===== SLIDE 9: Zusammenfassung =====
slide = add_content_slide(prs, "Zusammenfassung & Haupterkenntnisse")

add_text_box(slide, 0.5, 1.2, 9, 0.45, "3 Haupterkenntnisse aus TEIL A:",
             font_size=15, bold=True, color=DARK_BLUE)

insights_summary = [
    ("1. Monte Carlo ist zu langsam",
     "Updates nur am Episode-Ende → braucht ~15k Episodes → finale Performance: -123 (konvergiert nicht!)"),

    ("2. TD-Methoden sind revolutionär",
     "Online Updates nach jedem Schritt → SARSA & Q-Learning konvergieren bei ~5-10k Episodes"),

    ("3. Q-Learning ist aggressiv aber optimal",
     "Max-Bootstrap → schnelleres Lernen (40-50% speedup), am Ende ähnlich wie SARSA")
]

top = 1.8
for i, (title, desc) in enumerate(insights_summary):
    add_text_box(slide, 0.7, top, 8.8, 0.3, title,
                 font_size=12, bold=True, color=[RED, LIGHT_BLUE, GREEN][i])
    add_text_box(slide, 1.0, top + 0.35, 8.5, 0.45, desc,
                 font_size=11, color=DARK_GRAY)
    top += 1.05

add_text_box(slide, 0.5, 5.9, 9, 1.4,
             "🔑 DAS TD-FEHLER PARADOX:\n" +
             "Q-Learning hat GRÖSSERE TD-Fehler als SARSA, aber lernt SCHNELLER und erreicht bessere Performance!\n" +
             "→ TD-Fehler ist ein Lern-Signal, KEINE Qualitäts-Metrik!",
             font_size=12, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# ===== SLIDE 10: Details & Statistik =====
slide = add_content_slide(prs, "Detaillierte Statistik: Return und Episode Length")

stat_text = """
RETURN ANALYSE (Greedy Evaluation, Mean ± Std über 5 Seeds):

MC:         -123.23 ± 18.67  (Range: -147.45 bis -97.84)  ← Sehr schlecht, keine Konvergenz
SARSA:        7.92 ± 0.21    (Range: 7.73 bis 8.28)       ← Gut, sehr stabil
Q-Learning:   7.96 ± 0.22    (Range: 7.64 bis 8.28)       ← Leicht besser, stabil


EPISODE LENGTH ANALYSE (Durchschnitt über 300 Test-Episodes):

MC:         130.94 ± 16.78   (Range: 108.13 bis 152.70)   ← Lange Episodes, ineffizient
SARSA:       13.08 ± 0.21    (Range: 12.72 bis 13.27)     ← Kurz, effizient
Q-Learning:  13.04 ± 0.22    (Range: 12.72 bis 13.36)     ← Kurz, effizient


KONVERGENZ GESCHWINDIGKEIT (aus Learning Curves):

MC:         ~15.000 Episodes (aber erreicht nicht mal 0, bleibt bei -80)
SARSA:      ~8-10.000 Episodes (erreicht Return ~5-8)
Q-Learning: ~5-6.000 Episodes (erreicht Return ~5-8)  ← 40-50% schneller als SARSA!


VARIABILITÄT ÜBER SEEDS:

MC:         Hohe Streuung (±19) - Ergebnisse hängen stark vom Seed ab
SARSA:      Sehr stabil (±0.21) - kaum Unterschied zwischen Seeds
Q-Learning: Sehr stabil (±0.22) - kaum Unterschied zwischen Seeds
"""

add_text_box(slide, 0.3, 1.2, 9.4, 6.0, stat_text, font_size=9.5, color=DARK_GRAY)

# ============================================================================
# SPEICHERN
# ============================================================================

output_path = os.path.join(OUTPUT_DIR, "TEIL_A_Presentation_Final.pptx")
prs.save(output_path)

print(f"\n✓ Finale PowerPoint Präsentation erstellt:")
print(f"  {output_path}")
print(f"\n  Inhaltsverzeichnis:")
print(f"    Slide 0:  Title Slide")
print(f"    Slide 1:  Überblick & Aufgabenstellung")
print(f"    Slide 2:  Learning Curves (detailliert)")
print(f"    Slide 3:  Greedy Evaluation Bar Charts (detailliert)")
print(f"    Slide 4:  Resultate Tabelle & Insights")
print(f"    Slide 5:  Interpretation Q1: MC negativ?")
print(f"    Slide 6:  Interpretation Q2: TD schneller?")
print(f"    Slide 7:  Interpretation Q3: QL aggressiv?")
print(f"    Slide 8:  Interpretation Q4: SARSA & QL ähnlich?")
print(f"    Slide 9:  Zusammenfassung & Haupterkenntnisse")
print(f"    Slide 10: Detaillierte Statistik")
print(f"\n  Features:")
print(f"    ✓ Hochauflösende Grafiken (200 DPI)")
print(f"    ✓ Detaillierte Min/Max Streuung sichtbar")
print(f"    ✓ Alle 4 Interpretationsfragen beantwortet")
print(f"    ✓ Tabellen mit genauen Zahlen")
print(f"    ✓ Professionelles Design")
