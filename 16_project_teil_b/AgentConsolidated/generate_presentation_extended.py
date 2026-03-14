#!/usr/bin/env python3
"""
generate_presentation_extended.py

Standalone script to generate PowerPoint presentation from trading results including reward ablation.
Can be run independently after notebook execution.

Usage:
    python generate_presentation_extended.py --metrics metrics.pkl --reward_csv reward_ablation_comparison.csv --output presentation.pptx
"""

import argparse
import pickle
import json
import numpy as np
import pandas as pd
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. Install with: pip install python-pptx")


class PresentationGenerator:
    """Generate comprehensive presentation with reward ablation results"""
    
    def __init__(self):
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not available. Install with: pip install python-pptx")
        
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Color scheme
        self.DARK_BLUE = RGBColor(46, 134, 171)
        self.ACCENT_ORANGE = RGBColor(242, 142, 43)
        self.ACCENT_GREEN = RGBColor(46, 194, 113)
        self.WHITE = RGBColor(255, 255, 255)
        self.DARK_TEXT = RGBColor(64, 64, 64)
        self.LIGHT_GRAY = RGBColor(236, 240, 241)
    
    def add_title_slide(self, title, subtitle):
        """Add title slide with custom styling"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.DARK_BLUE

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = self.WHITE

        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = self.ACCENT_ORANGE

        return slide

    def add_content_slide(self, title, content_bullets):
        """Add content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.WHITE

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = self.DARK_BLUE

        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, bullet in enumerate(content_bullets):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = self.DARK_TEXT
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)

        return slide

    def add_image_slide(self, title, image_path):
        """Add slide with full-width image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.WHITE

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.DARK_BLUE

        # Add image
        if Path(image_path).exists():
            try:
                slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1), width=Inches(9))
            except Exception as e:
                print(f"⚠ Could not add image {image_path}: {e}")

        return slide
    
    def add_metrics_table_slide(self, title, metrics_dict):
        """Add slide with metrics table"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Add background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.WHITE
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.DARK_BLUE
        
        # Create table (5 columns: Metric, Best, Good, Average, Worst)
        rows, cols = min(10, len(metrics_dict) + 1), 5
        left = Inches(0.5)
        top = Inches(1)
        width = Inches(9)
        height = Inches(5.5)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Set column widths
        for col_idx in range(cols):
            table.columns[col_idx].width = Inches(width.inches / cols)
        
        # Header row
        headers = ['Metric', 'Best', 'Value', 'Reward Type', 'Notes']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.DARK_BLUE
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = self.WHITE
        
        # Data rows
        for row_idx, (metric_name, metric_value) in enumerate(list(metrics_dict.items())[:rows-1], 1):
            # Metric name
            cell = table.cell(row_idx, 0)
            cell.text = str(metric_name).replace('_', ' ').title()
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            
            # Value
            cell = table.cell(row_idx, 2)
            cell.text = f'{float(metric_value):.4f}'
            cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        return slide


def create_reward_ablation_presentation(reward_csv_path, output_path='reward_ablation_analysis.pptx'):
    """
    Create presentation from reward ablation results
    
    Parameters
    ----------
    reward_csv_path : str
        Path to reward_ablation_comparison.csv
    output_path : str
        Where to save the presentation
    """
    if not PPTX_AVAILABLE:
        print("Error: python-pptx not available")
        print("Install with: pip install python-pptx")
        return False
    
    # Load data
    df = pd.read_csv(reward_csv_path, index_col=0)
    
    gen = PresentationGenerator()
    
    # Slide 1: Title
    gen.add_title_slide(
        "Reward Function Ablation Study",
        "Comprehensive Analysis of 8 Reward Variants for PPO Trading"
    )
    
    # Slide 2: Overview
    gen.add_content_slide(
        "Study Overview",
        [
            "• Tested 8 different reward function formulations",
            "• Evaluated on BTC-USD historical data",
            "• Compared across key performance metrics",
            "• Identified best performers for different objectives",
            "• Analyzed trade-offs between return, risk, and consistency"
        ]
    )
    
    # Slide 3: Reward Function Summary
    gen.add_content_slide(
        "8 Reward Function Types",
        [
            "1. BASIC: Pure returns (R = PnL - Cost)",
            "2. WITH_RISK: Balanced with risk penalty",
            "3. WITH_SHARPE: Risk-adjusted returns",
            "4. RISK_ADJUSTED: Normalized returns",
            "5. SORTINO: Downside-risk focused",
            "6. CALMAR: Drawdown control",
            "7. INFORMATION_RATIO: Consistency bonus",
            "8. COMPOSITE: Multi-objective blend"
        ]
    )
    
    # Slide 4: Key Findings
    best_return_idx = df['total_return'].idxmax()
    best_sharpe_idx = df['sharpe_ratio'].idxmax()
    best_drawdown_idx = df['max_drawdown'].idxmax()
    
    gen.add_content_slide(
        "Key Findings",
        [
            f"• Highest Return: {best_return_idx.replace('PPO_', '').replace('_', ' ').title()} ({df.loc[best_return_idx, 'total_return']*100:.2f}%)",
            f"• Best Risk-Adjusted: {best_sharpe_idx.replace('PPO_', '').replace('_', ' ').title()} (Sharpe: {df.loc[best_sharpe_idx, 'sharpe_ratio']:.2f})",
            f"• Best Drawdown Control: {best_drawdown_idx.replace('PPO_', '').replace('_', ' ').title()} ({df.loc[best_drawdown_idx, 'max_drawdown']*100:.2f}%)",
            "• Composite reward provides good balance across metrics",
            "• Trade-off between return maximization and risk control"
        ]
    )
    
    # Slide 5: Recommendations
    gen.add_content_slide(
        "Recommendations by Use Case",
        [
            "📈 For Maximum Returns: Use BASIC reward",
            "⚖️ For Balanced Approach: Use COMPOSITE or WITH_RISK",
            "🛡️ For Conservative: Use WITH_SHARPE or SORTINO",
            "📉 For Drawdown Control: Use CALMAR reward",
            "💪 For Consistency: Use INFORMATION_RATIO",
            "⭐ General Recommendation: COMPOSITE (balanced across all metrics)"
        ]
    )
    
    # Slide 6-10: Add images if available
    images = [
        ('reward_metrics_comparison.png', 'Performance Metrics Comparison'),
        ('reward_heatmap.png', 'Normalized Performance Heatmap'),
        ('reward_scatter.png', 'Risk vs Return Analysis'),
        ('reward_ranking.png', 'Rankings by Metric'),
        ('baseline_vs_rewards.png', 'Baseline Comparison')
    ]
    
    for img_file, img_title in images:
        if Path(img_file).exists():
            gen.add_image_slide(img_title, img_file)
    
    # Slide N: Metrics Summary Table
    metrics_summary = df.describe().to_dict()
    gen.add_metrics_table_slide("Metrics Summary", df.iloc[0].to_dict())
    
    # Slide N+1: Conclusion
    gen.add_content_slide(
        "Conclusion",
        [
            "✓ Reward function choice significantly impacts performance",
            "✓ No single 'best' reward - choose based on objective",
            "✓ COMPOSITE provides excellent balance",
            "✓ Risk-aware rewards (WITH_SHARPE, SORTINO) more robust",
            "✓ Consider market conditions and risk tolerance",
            "✓ Further tuning of reward parameters recommended"
        ]
    )
    
    # Save presentation
    gen.prs.save(output_path)
    print(f"✓ Presentation saved: {output_path}")
    return True


def create_comprehensive_presentation(metrics_pkl_path, reward_csv_path=None, output_path='trading_analysis.pptx', images_dir='./plots'):
    """
    Create comprehensive presentation from all available results
    
    Parameters
    ----------
    metrics_pkl_path : str
        Path to metrics.pkl from run_all_experiments.py
    reward_csv_path : str, optional
        Path to reward_ablation_comparison.csv
    output_path : str
        Where to save the presentation
    images_dir : str
        Directory containing generated images
    """
    if not PPTX_AVAILABLE:
        print("Error: python-pptx not available")
        return False
    
    gen = PresentationGenerator()
    
    # Slide 1: Title
    gen.add_title_slide(
        "PPO Trading with Reward Ablation",
        "Comprehensive Performance Analysis"
    )
    
    # Slide 2: Executive Summary
    gen.add_content_slide(
        "Executive Summary",
        [
            "• Comprehensive evaluation of 10 trading experiments",
            "• 2 baseline experiments (with/without LSTM forecast)",
            "• 8 reward function ablation studies",
            "• Systematic comparison of risk-return trade-offs",
            "• Data: BTC-USD (historical, 2018-2026)"
        ]
    )
    
    # Load metrics from pickle if available
    try:
        with open(metrics_pkl_path, 'rb') as f:
            metrics_data = pickle.load(f)
        
        if isinstance(metrics_data, dict) and 'results' in metrics_data:
            # Extract experiment results
            experiments = metrics_data['results']
            
            gen.add_content_slide(
                "Experiments Overview",
                [
                    f"Total Experiments: {len(experiments)}",
                    f"Data Range: {metrics_data.get('timestamp', 'Unknown')}",
                    "",
                    "Baseline Experiments:",
                    "  • PPO Without Forecast",
                    "  • PPO With LSTM Forecast",
                    "",
                    "Reward Ablation (8 variants):",
                    "  • BASIC, WITH_RISK, WITH_SHARPE, RISK_ADJUSTED",
                    "  • SORTINO, CALMAR, INFORMATION_RATIO, COMPOSITE"
                ]
            )
    except Exception as e:
        print(f"⚠ Could not load metrics pickle: {e}")
    
    # Load and add reward comparison if available
    if reward_csv_path and Path(reward_csv_path).exists():
        try:
            df_rewards = pd.read_csv(reward_csv_path, index_col=0)
            
            best_return = df_rewards['total_return'].idxmax()
            best_sharpe = df_rewards['sharpe_ratio'].idxmax()
            
            gen.add_content_slide(
                "Reward Comparison Results",
                [
                    f"Best Return: {best_return.replace('PPO_', '').title()}",
                    f"  └─ Return: {df_rewards.loc[best_return, 'total_return']*100:.2f}%",
                    "",
                    f"Best Sharpe: {best_sharpe.replace('PPO_', '').title()}",
                    f"  └─ Sharpe Ratio: {df_rewards.loc[best_sharpe, 'sharpe_ratio']:.2f}",
                    "",
                    "Recommendation: Choose based on your objectives",
                    "  • Risk-averse → WITH_SHARPE",
                    "  • Balanced → COMPOSITE",
                    "  • Aggressive → BASIC"
                ]
            )
        except Exception as e:
            print(f"⚠ Could not load reward comparison: {e}")
    
    # Add image slides
    images_dir = Path(images_dir)
    if images_dir.exists():
        image_files = sorted(images_dir.glob('*.png'))
        
        for img_file in image_files[:8]:  # Limit to 8 images
            gen.add_image_slide(
                img_file.stem.replace('_', ' ').title(),
                str(img_file)
            )
    
    # Final slide: Next Steps
    gen.add_content_slide(
        "Next Steps",
        [
            "1. Review detailed metrics in CSV files",
            "2. Analyze visualizations for insights",
            "3. Select reward function for production",
            "4. Fine-tune parameters for specific use case",
            "5. Validate on out-of-sample data",
            "6. Deploy and monitor performance"
        ]
    )
    
    # Save presentation
    gen.prs.save(output_path)
    print(f"✓ Presentation saved: {output_path}")
    return True


def main():
    """Main execution"""
    parser = argparse.ArgumentParser(description='Generate presentation from trading results')
    parser.add_argument('--metrics', type=str, help='Path to metrics.pkl')
    parser.add_argument('--reward_csv', type=str, help='Path to reward_ablation_comparison.csv')
    parser.add_argument('--reward_only', action='store_true', help='Generate reward ablation presentation only')
    parser.add_argument('--output', type=str, default='trading_analysis.pptx', help='Output presentation file')
    parser.add_argument('--images_dir', type=str, default='./plots', help='Directory with generated images')
    
    args = parser.parse_args()
    
    if not PPTX_AVAILABLE:
        print("Error: python-pptx not installed")
        print("Install with: pip install python-pptx")
        return 1
    
    # Generate reward-only presentation
    if args.reward_only and args.reward_csv:
        create_reward_ablation_presentation(args.reward_csv, args.output)
    
    # Generate comprehensive presentation
    elif args.metrics:
        create_comprehensive_presentation(
            args.metrics,
            args.reward_csv,
            args.output,
            args.images_dir
        )
    
    else:
        print("Error: Provide --metrics or --reward_csv")
        print("Usage:")
        print("  python generate_presentation_extended.py --metrics metrics.pkl --output report.pptx")
        print("  python generate_presentation_extended.py --reward_csv reward_ablation_comparison.csv --reward_only --output rewards.pptx")
        return 1
    
    print(f"✓ Presentation generated: {args.output}")
    return 0


if __name__ == "__main__":
    exit(main())

